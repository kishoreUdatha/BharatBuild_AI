"""
PowerPoint Presentation Generator for Academic Projects
Generates professional PPT presentations with images from project data
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from typing import Dict, List, Optional
from datetime import datetime
import os

from app.core.logging_config import logger


class AcademicPPTGenerator:
    """
    Generate professional PowerPoint presentations for student projects

    Features:
    - Professional slide designs with NO content overlap
    - Title slide, agenda, content, conclusion
    - Images and diagrams from project
    - Consistent branding
    - 20-25 slides
    - Maximum 5 bullet points per slide to prevent overlap
    """

    # Maximum bullet points per slide to prevent overlap
    MAX_BULLETS_PER_SLIDE = 5

    def __init__(self):
        # Professional color scheme
        self.primary_color = RGBColor(44, 62, 80)      # Dark blue-gray
        self.secondary_color = RGBColor(52, 73, 94)    # Medium blue-gray
        self.accent_color = RGBColor(41, 128, 185)     # Bright blue
        self.text_color = RGBColor(51, 51, 51)         # Dark gray
        self.light_gray = RGBColor(236, 240, 241)      # Light background
        self.success_color = RGBColor(46, 204, 113)    # Green

    def generate_project_presentation(
        self,
        ppt_data: Dict,
        output_path: str,
        project_info: Optional[Dict] = None,
        diagrams_dir: Optional[str] = None
    ) -> bool:
        """
        Generate complete PowerPoint presentation with images

        Args:
            ppt_data: Presentation content from Document Generator Agent
            output_path: Path to save .pptx file
            project_info: Project metadata (title, authors, etc.)
            diagrams_dir: Directory containing UML diagrams for the project

        Returns:
            bool: True if successful
        """
        try:
            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(10)  # 16:9 aspect ratio
            prs.slide_height = Inches(5.625)

            # Get slide data
            slides_data = ppt_data.get('slides', [])
            project_title = ppt_data.get('project_title', 'Project Presentation')

            # Find available diagrams
            available_diagrams = self._find_diagrams(diagrams_dir) if diagrams_dir else {}

            # Generate slides
            for slide_data in slides_data:
                slide_number = slide_data.get('slide_number', 0)
                title = slide_data.get('title', '')
                content = slide_data.get('content', {})
                content_type = content.get('type', 'bullets')

                if slide_number == 1:
                    self._create_title_slide(prs, title, content, project_info)
                elif content_type == 'agenda':
                    self._create_agenda_slide(prs, title, content)
                elif content_type == 'bullets':
                    self._create_bullet_slide_with_image(prs, title, content, available_diagrams)
                elif content_type == 'two_column':
                    self._create_two_column_slide(prs, title, content)
                elif content_type == 'architecture':
                    self._create_architecture_slide(prs, title, content, available_diagrams)
                elif content_type == 'features':
                    self._create_features_slide(prs, title, content)
                elif content_type == 'tech_stack':
                    self._create_tech_stack_slide(prs, title, content)
                elif content_type == 'database_schema':
                    self._create_database_slide(prs, title, content, available_diagrams)
                elif content_type == 'testing_results':
                    self._create_testing_slide(prs, title, content)
                elif content_type == 'conclusion':
                    self._create_conclusion_slide(prs, title, content)
                elif content_type == 'future_enhancements':
                    self._create_future_scope_slide(prs, title, content)
                elif content_type == 'thank_you':
                    self._create_thank_you_slide(prs, title, content)
                else:
                    # Default bullet slide
                    self._create_bullet_slide_with_image(prs, title, content, available_diagrams)

            # Save presentation
            prs.save(output_path)

            logger.info(f"Generated PowerPoint presentation: {output_path}")
            return True

        except Exception as e:
            logger.error(f"Error generating PowerPoint: {e}", exc_info=True)
            return False

    def _find_diagrams(self, diagrams_dir: str) -> Dict[str, str]:
        """Find available UML diagrams in directory"""
        diagrams = {}
        if not diagrams_dir or not os.path.exists(diagrams_dir):
            return diagrams

        diagram_types = {
            'use_case': ['use_case', 'usecase'],
            'class': ['class_diagram', 'class'],
            'sequence': ['sequence', 'seq'],
            'activity': ['activity', 'flow'],
            'er': ['er_diagram', 'erd', 'entity'],
            'dfd': ['dfd', 'data_flow'],
            'architecture': ['architecture', 'system_arch', 'component']
        }

        for filename in os.listdir(diagrams_dir):
            if filename.lower().endswith('.png'):
                filepath = os.path.join(diagrams_dir, filename)
                filename_lower = filename.lower()
                for diagram_type, keywords in diagram_types.items():
                    for keyword in keywords:
                        if keyword in filename_lower:
                            diagrams[diagram_type] = filepath
                            break

        return diagrams

    def _truncate_points(self, points: List[str], max_points: int = None) -> List[str]:
        """Truncate bullet points to prevent overflow"""
        max_pts = max_points or self.MAX_BULLETS_PER_SLIDE
        if len(points) <= max_pts:
            return points
        # Take first max-1 and add "...and more"
        truncated = points[:max_pts-1]
        remaining = len(points) - (max_pts - 1)
        truncated.append(f"...and {remaining} more items")
        return truncated

    def _create_title_slide(self, prs: Presentation, title: str, content: Dict, project_info: Optional[Dict]):
        """Create title slide"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.primary_color

        # Project title
        project_title = content.get('project_name', title)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.2))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = project_title
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle = content.get('subtitle', 'Academic Project Presentation')
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.color.rgb = self.light_gray
        subtitle_para.alignment = PP_ALIGN.CENTER

        # Student/Author info
        if project_info and 'authors' in project_info:
            authors_text = ", ".join(project_info['authors'][:3])  # Max 3 authors
        else:
            authors_text = content.get('presented_by', 'Student Name')

        author_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.6))
        author_frame = author_box.text_frame
        author_para = author_frame.paragraphs[0]
        author_para.text = f"Presented By: {authors_text}"
        author_para.font.size = Pt(16)
        author_para.font.color.rgb = RGBColor(255, 255, 255)
        author_para.alignment = PP_ALIGN.CENTER

        # Date
        date_text = content.get('date', datetime.utcnow().strftime('%B %d, %Y'))
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(9), Inches(0.3))
        date_frame = date_box.text_frame
        date_para = date_frame.paragraphs[0]
        date_para.text = date_text
        date_para.font.size = Pt(14)
        date_para.font.color.rgb = self.light_gray
        date_para.alignment = PP_ALIGN.CENTER

    def _create_agenda_slide(self, prs: Presentation, title: str, content: Dict):
        """Create agenda/outline slide"""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Title bar
        self._add_title_bar(slide, title)

        # Content
        items = content.get('items', [])
        items = self._truncate_points(items, 8)  # Allow 8 for agenda

        left = Inches(1.5)
        top = Inches(1.8)
        width = Inches(7)
        height = Inches(3.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        for i, item in enumerate(items, 1):
            p = text_frame.add_paragraph() if i > 1 else text_frame.paragraphs[0]
            p.text = f"{i}. {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.text_color
            p.space_before = Pt(8)
            p.level = 0

    def _create_bullet_slide_with_image(self, prs: Presentation, title: str, content: Dict, diagrams: Dict):
        """Create bullet point slide with optional image"""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Title bar
        self._add_title_bar(slide, title)

        # Determine if we should add an image
        title_lower = title.lower()
        image_path = None

        # Match diagram based on slide title
        if 'architecture' in title_lower or 'system design' in title_lower:
            image_path = diagrams.get('architecture')
        elif 'use case' in title_lower:
            image_path = diagrams.get('use_case')
        elif 'class' in title_lower:
            image_path = diagrams.get('class')
        elif 'sequence' in title_lower:
            image_path = diagrams.get('sequence')
        elif 'activity' in title_lower or 'flow' in title_lower:
            image_path = diagrams.get('activity')
        elif 'database' in title_lower or 'er' in title_lower or 'entity' in title_lower:
            image_path = diagrams.get('er')
        elif 'dfd' in title_lower or 'data flow' in title_lower:
            image_path = diagrams.get('dfd')

        # Get bullet points
        points = content.get('points', [])
        points = self._truncate_points(points)

        if image_path and os.path.exists(image_path):
            # Two-column layout: bullets on left, image on right
            # Left column - bullets
            left = Inches(0.5)
            top = Inches(1.8)
            width = Inches(4.5)
            height = Inches(3.3)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            for i, point in enumerate(points[:4]):  # Max 4 points with image
                p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
                clean_point = self._clean_text(point)
                p.text = f"• {clean_point}"
                p.font.size = Pt(14)
                p.font.color.rgb = self.text_color
                p.space_before = Pt(6)

            # Right side - image
            try:
                slide.shapes.add_picture(
                    image_path,
                    Inches(5.2),
                    Inches(1.6),
                    width=Inches(4.3),
                    height=Inches(3.5)
                )
            except Exception as e:
                logger.warning(f"Could not add image to slide: {e}")
        else:
            # Full width bullets
            left = Inches(1)
            top = Inches(1.8)
            width = Inches(8)
            height = Inches(3.3)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            for i, point in enumerate(points):
                p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
                clean_point = self._clean_text(point)
                p.text = f"• {clean_point}"
                p.font.size = Pt(16)
                p.font.color.rgb = self.text_color
                p.space_before = Pt(8)

    def _add_title_bar(self, slide, title: str):
        """Add consistent title bar to slide"""
        # Title background bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(10),
            Inches(1.3)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.primary_color
        title_bar.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

    def _clean_text(self, text: str) -> str:
        """Remove emojis and clean text"""
        # Remove common emojis
        emoji_pattern = re.compile("["
            u"\U0001F600-\U0001F64F"  # emoticons
            u"\U0001F300-\U0001F5FF"  # symbols & pictographs
            u"\U0001F680-\U0001F6FF"  # transport & map symbols
            u"\U0001F700-\U0001F77F"  # alchemical symbols
            u"\U0001F780-\U0001F7FF"  # Geometric Shapes Extended
            u"\U0001F800-\U0001F8FF"  # Supplemental Arrows-C
            u"\U0001F900-\U0001F9FF"  # Supplemental Symbols and Pictographs
            u"\U0001FA00-\U0001FA6F"  # Chess Symbols
            u"\U0001FA70-\U0001FAFF"  # Symbols and Pictographs Extended-A
            u"\U00002702-\U000027B0"  # Dingbats
            u"\U0001F1E0-\U0001F1FF"  # flags (iOS)
            "]+", flags=re.UNICODE)

        text = emoji_pattern.sub('', text)
        text = text.lstrip('✅❌🎯📊💡🔧⚡🚀📝📌⭐🔥💻📁')
        return text.strip()

    def _create_two_column_slide(self, prs: Presentation, title: str, content: Dict):
        """Create two-column slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title bar
        self._add_title_bar(slide, title)

        # Left column
        left_content = content.get('left', {})
        left_title = left_content.get('title', 'Left')
        left_points = self._truncate_points(left_content.get('points', []), 4)

        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.3), Inches(3.3))
        left_frame = left_box.text_frame

        # Left title
        p = left_frame.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.accent_color

        # Left points
        for point in left_points:
            p = left_frame.add_paragraph()
            p.text = f"• {self._clean_text(point)}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.text_color
            p.space_before = Pt(6)

        # Right column
        right_content = content.get('right', {})
        right_title = right_content.get('title', 'Right')
        right_points = self._truncate_points(right_content.get('points', []), 4)

        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.3), Inches(3.3))
        right_frame = right_box.text_frame

        # Right title
        p = right_frame.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.accent_color

        # Right points
        for point in right_points:
            p = right_frame.add_paragraph()
            p.text = f"• {self._clean_text(point)}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.text_color
            p.space_before = Pt(6)

    def _create_architecture_slide(self, prs: Presentation, title: str, content: Dict, diagrams: Dict):
        """Create system architecture slide with diagram"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title bar
        self._add_title_bar(slide, title)

        # Check for architecture diagram
        arch_image = diagrams.get('architecture')

        if arch_image and os.path.exists(arch_image):
            # Add diagram
            try:
                slide.shapes.add_picture(
                    arch_image,
                    Inches(1.5),
                    Inches(1.6),
                    width=Inches(7),
                    height=Inches(3.7)
                )
                return
            except Exception as e:
                logger.warning(f"Could not add architecture diagram: {e}")

        # Fallback: Show layers as boxes
        layers = content.get('layers', [])
        if not layers:
            layers = ['Presentation Layer', 'Business Logic Layer', 'Data Access Layer', 'Database Layer']

        layers = layers[:5]  # Max 5 layers

        left = Inches(2)
        top = Inches(1.8)
        width = Inches(6)
        layer_height = Inches(0.7)

        colors = [
            RGBColor(52, 152, 219),   # Blue
            RGBColor(46, 204, 113),   # Green
            RGBColor(155, 89, 182),   # Purple
            RGBColor(241, 196, 15),   # Yellow
            RGBColor(230, 126, 34),   # Orange
        ]

        for i, layer in enumerate(layers):
            # Create colored box for each layer
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left,
                top + (i * (layer_height + Inches(0.1))),
                width,
                layer_height
            )

            shape.fill.solid()
            shape.fill.fore_color.rgb = colors[i % len(colors)]
            shape.line.fill.background()

            # Layer text
            text_frame = shape.text_frame
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            p = text_frame.paragraphs[0]
            p.text = self._clean_text(layer)
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

    def _create_features_slide(self, prs: Presentation, title: str, content: Dict):
        """Create features slide with icons"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title bar
        self._add_title_bar(slide, title)

        points = content.get('points', [])
        points = self._truncate_points(points, 6)

        # Create a grid of feature boxes
        cols = 2
        rows = 3
        box_width = Inches(4.2)
        box_height = Inches(0.9)
        start_left = Inches(0.7)
        start_top = Inches(1.7)
        h_gap = Inches(0.4)
        v_gap = Inches(0.15)

        for i, point in enumerate(points):
            row = i // cols
            col = i % cols

            left = start_left + col * (box_width + h_gap)
            top = start_top + row * (box_height + v_gap)

            # Feature box
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left,
                top,
                box_width,
                box_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.light_gray
            shape.line.color.rgb = self.accent_color

            # Feature text
            text_frame = shape.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = f"✓ {self._clean_text(point)}"
            p.font.size = Pt(13)
            p.font.color.rgb = self.text_color
            p.alignment = PP_ALIGN.LEFT

    def _create_tech_stack_slide(self, prs: Presentation, title: str, content: Dict):
        """Create technology stack slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title bar
        self._add_title_bar(slide, title)

        # Technologies by category
        categories = content.get('categories', {})

        left = Inches(1)
        top = Inches(1.8)
        width = Inches(8)

        textbox = slide.shapes.add_textbox(left, top, width, Inches(3.3))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        cat_count = 0
        for category, tech_list in categories.items():
            if cat_count >= 4:  # Max 4 categories
                break

            # Category header
            p = text_frame.add_paragraph() if cat_count > 0 else text_frame.paragraphs[0]
            p.text = f"{category}:"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.accent_color
            p.space_before = Pt(10)

            # Technologies
            if isinstance(tech_list, list):
                techs = tech_list[:4]  # Max 4 techs per category
                p = text_frame.add_paragraph()
                p.text = "   " + ", ".join(techs)
                p.font.size = Pt(14)
                p.font.color.rgb = self.text_color
            else:
                p = text_frame.add_paragraph()
                p.text = f"   {tech_list}"
                p.font.size = Pt(14)
                p.font.color.rgb = self.text_color

            cat_count += 1

    def _create_database_slide(self, prs: Presentation, title: str, content: Dict, diagrams: Dict):
        """Create database schema slide with ER diagram"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title bar
        self._add_title_bar(slide, title)

        # Check for ER diagram
        er_image = diagrams.get('er')

        if er_image and os.path.exists(er_image):
            try:
                slide.shapes.add_picture(
                    er_image,
                    Inches(1.5),
                    Inches(1.6),
                    width=Inches(7),
                    height=Inches(3.7)
                )
                return
            except Exception as e:
                logger.warning(f"Could not add ER diagram: {e}")

        # Fallback: bullet points
        tables = content.get('tables', content.get('points', []))
        tables = self._truncate_points(tables)

        left = Inches(1)
        top = Inches(1.8)
        width = Inches(8)

        textbox = slide.shapes.add_textbox(left, top, width, Inches(3.3))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        for i, table in enumerate(tables):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = f"• {self._clean_text(table)}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.text_color
            p.space_before = Pt(8)

    def _create_testing_slide(self, prs: Presentation, title: str, content: Dict):
        """Create testing results slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title bar
        self._add_title_bar(slide, title)

        points = content.get('points', [])
        points = self._truncate_points(points)

        left = Inches(1)
        top = Inches(1.8)
        width = Inches(8)

        textbox = slide.shapes.add_textbox(left, top, width, Inches(3.3))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        for i, point in enumerate(points):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            clean_point = self._clean_text(point)
            # Add checkmark for test results
            p.text = f"✓ {clean_point}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.success_color
            p.space_before = Pt(8)

    def _create_conclusion_slide(self, prs: Presentation, title: str, content: Dict):
        """Create conclusion slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title bar
        self._add_title_bar(slide, title)

        points = content.get('points', [])
        points = self._truncate_points(points)

        left = Inches(1)
        top = Inches(1.8)
        width = Inches(8)

        textbox = slide.shapes.add_textbox(left, top, width, Inches(3.3))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        for i, point in enumerate(points):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = f"• {self._clean_text(point)}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.text_color
            p.space_before = Pt(8)

    def _create_future_scope_slide(self, prs: Presentation, title: str, content: Dict):
        """Create future scope slide"""
        enhancements = content.get('enhancements', content.get('points', []))
        enhancements = self._truncate_points(enhancements)
        self._create_bullet_slide_with_image(
            prs, title, {'points': enhancements}, {}
        )

    def _create_thank_you_slide(self, prs: Presentation, title: str, content: Dict):
        """Create thank you slide"""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.primary_color

        # Thank you text
        thank_you_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(0.8))
        thank_you_frame = thank_you_box.text_frame
        thank_you_para = thank_you_frame.paragraphs[0]
        thank_you_para.text = "Thank You!"
        thank_you_para.font.size = Pt(48)
        thank_you_para.font.bold = True
        thank_you_para.font.color.rgb = RGBColor(255, 255, 255)
        thank_you_para.alignment = PP_ALIGN.CENTER

        # Questions text
        message = content.get('message', 'Questions?')
        questions_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.5))
        questions_frame = questions_box.text_frame
        questions_para = questions_frame.paragraphs[0]
        questions_para.text = message
        questions_para.font.size = Pt(28)
        questions_para.font.color.rgb = self.light_gray
        questions_para.alignment = PP_ALIGN.CENTER

        # Contact info
        contact = content.get('contact', '')
        if contact:
            contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.4))
            contact_frame = contact_box.text_frame
            contact_para = contact_frame.paragraphs[0]
            contact_para.text = contact
            contact_para.font.size = Pt(16)
            contact_para.font.color.rgb = RGBColor(255, 255, 255)
            contact_para.alignment = PP_ALIGN.CENTER

        # Powered by BharatBuild AI
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.3))
        footer_frame = footer_box.text_frame
        footer_para = footer_frame.paragraphs[0]
        footer_para.text = "Generated by BharatBuild AI"
        footer_para.font.size = Pt(12)
        footer_para.font.italic = True
        footer_para.font.color.rgb = self.light_gray
        footer_para.alignment = PP_ALIGN.CENTER


# Singleton instance
ppt_generator = AcademicPPTGenerator()
