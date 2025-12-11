"""
POWERPOINT GENERATOR V2
=======================
Creates beautiful 20-25 slide presentations with professional design.

Features:
- Modern slide design with color scheme
- Smart layouts for different content types
- Bullet points with icons
- Code snippets with syntax highlighting look
- Architecture diagrams as shapes
- Screenshots placeholders
- Animations suggestions
"""

from typing import Dict, List, Optional, Any
from datetime import datetime
import os
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

from app.core.logging_config import logger
from app.core.config import settings


class PPTGeneratorV2:
    """
    Professional PowerPoint Generator with modern design

    Creates beautiful presentations with:
    - Custom color scheme
    - Professional layouts
    - Smart content formatting
    """

    # Color Scheme (Modern Dark Blue Theme)
    PRIMARY_COLOR = RGBColor(0, 51, 102)      # Dark blue
    SECONDARY_COLOR = RGBColor(0, 102, 153)   # Teal
    ACCENT_COLOR = RGBColor(255, 153, 0)      # Orange
    TEXT_COLOR = RGBColor(51, 51, 51)         # Dark gray
    LIGHT_BG = RGBColor(240, 244, 248)        # Light gray-blue
    WHITE = RGBColor(255, 255, 255)

    # Slide dimensions (16:9)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    def __init__(self):
        self.prs = None

    async def create_presentation(
        self,
        sections: List[Dict],
        project_data: Dict,
        project_id: str = None,
        user_id: str = None
    ) -> str:
        """
        Create complete PowerPoint presentation.

        Args:
            sections: List of slide sections
            project_data: Project metadata
            project_id: Project ID for saving to project's docs folder
            user_id: User ID for isolation

        Returns:
            Path to generated presentation
        """
        # Store user_id for diagram generation
        self.user_id = user_id
        self.project_id = project_id

        try:
            # Create presentation
            self.prs = Presentation()
            self.prs.slide_width = self.SLIDE_WIDTH
            self.prs.slide_height = self.SLIDE_HEIGHT

            # Add slides for each section
            for section in sections:
                await self._add_section_slides(section, project_data)

            # Determine output directory - prefer project docs folder with user isolation
            if project_id:
                output_dir = settings.get_project_docs_dir(project_id, user_id)
            else:
                output_dir = settings.GENERATED_DIR / "presentations"
            output_dir.mkdir(parents=True, exist_ok=True)

            filename = f"{project_data.get('project_name', 'Presentation')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            file_path = output_dir / filename

            self.prs.save(str(file_path))

            logger.info(f"[PPTGenerator] Created presentation: {file_path}")

            return str(file_path)

        except Exception as e:
            logger.error(f"[PPTGenerator] Error: {e}", exc_info=True)
            raise

    async def _add_section_slides(self, section: Dict, project_data: Dict):
        """Add slides for a section"""
        section_type = section.get("type", "generate")
        section_id = section.get("section_id", "")
        content = section.get("content", {})

        if section_type == "template":
            if section_id == "title":
                self._add_title_slide(content, project_data)
            elif section_id == "thankyou":
                self._add_thankyou_slide(content, project_data)
        else:
            self._add_content_slide(section, project_data)

    def _add_title_slide(self, content: Dict, project_data: Dict):
        """Add title slide with professional design"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Background gradient effect (using shape)
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.PRIMARY_COLOR
        bg_shape.line.fill.background()

        # Accent bar at top
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, Inches(0.3)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.ACCENT_COLOR
        accent_bar.line.fill.background()

        # Project title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True

        p = title_frame.paragraphs[0]
        p.text = content.get("project_name", project_data.get("project_name", "Project Title"))
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2),
            Inches(12.333), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame

        p = subtitle_frame.paragraphs[0]
        p.text = content.get("subtitle", project_data.get("project_type", "Software Project"))
        p.font.size = Pt(24)
        p.font.color.rgb = self.LIGHT_BG
        p.alignment = PP_ALIGN.CENTER

        # Presented by
        presenter_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.5),
            Inches(12.333), Inches(1)
        )
        presenter_frame = presenter_box.text_frame

        p = presenter_frame.paragraphs[0]
        p.text = f"Presented by: {content.get('presented_by', project_data.get('author', 'Student Name'))}"
        p.font.size = Pt(18)
        p.font.color.rgb = self.LIGHT_BG
        p.alignment = PP_ALIGN.CENTER

        # Guide info
        p = presenter_frame.add_paragraph()
        p.text = f"Guide: {content.get('guide', project_data.get('guide', 'Guide Name'))}"
        p.font.size = Pt(16)
        p.font.color.rgb = self.LIGHT_BG
        p.alignment = PP_ALIGN.CENTER

    def _add_content_slide(self, section: Dict, project_data: Dict):
        """Add content slide with professional layout"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        title = section.get("title", "")
        content = section.get("content", {})

        # Header bar
        header_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, Inches(1.2)
        )
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = self.PRIMARY_COLOR
        header_bar.line.fill.background()

        # Accent line
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(1.2),
            self.SLIDE_WIDTH, Inches(0.05)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = self.ACCENT_COLOR
        accent_line.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame

        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

        # Content area
        if isinstance(content, dict):
            self._add_structured_content(slide, content)
        elif isinstance(content, str):
            self._add_text_content(slide, content)

    def _add_structured_content(self, slide, content: Dict):
        """Add structured content (subsections, bullets, etc.)"""
        # Main content box
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(12.333), Inches(5.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        # Get content text
        text = content.get("content", "")
        subsections = content.get("subsections", [])

        # Add main content
        if text:
            bullets = self._extract_bullets(text)
            for bullet in bullets[:6]:  # Max 6 bullets
                p = content_frame.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(20)
                p.font.color.rgb = self.TEXT_COLOR
                p.space_before = Pt(12)
                p.level = 0

        # Add subsection content
        for subsection in subsections[:3]:  # Max 3 subsections
            sub_title = subsection.get("title", "")
            sub_content = subsection.get("content", "")

            if sub_title:
                p = content_frame.add_paragraph()
                p.text = sub_title
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = self.SECONDARY_COLOR
                p.space_before = Pt(18)

            if sub_content:
                bullets = self._extract_bullets(sub_content)
                for bullet in bullets[:4]:  # Max 4 bullets per subsection
                    p = content_frame.add_paragraph()
                    p.text = f"  • {bullet}"
                    p.font.size = Pt(18)
                    p.font.color.rgb = self.TEXT_COLOR
                    p.space_before = Pt(6)

    def _add_text_content(self, slide, text: str):
        """Add text-only content"""
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(12.333), Inches(5.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        bullets = self._extract_bullets(text)
        for bullet in bullets[:6]:
            p = content_frame.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(20)
            p.font.color.rgb = self.TEXT_COLOR
            p.space_before = Pt(12)

    def _extract_bullets(self, text: str) -> List[str]:
        """Extract bullet points from text"""
        # Split by common delimiters
        if '\n' in text:
            lines = text.split('\n')
        elif '.' in text:
            lines = [s.strip() + '.' for s in text.split('.') if s.strip()]
        else:
            lines = [text]

        bullets = []
        for line in lines:
            line = line.strip()
            # Remove existing bullet markers
            line = line.lstrip('-•* ')
            if line and len(line) > 5:  # Skip very short lines
                # Truncate long lines
                if len(line) > 100:
                    line = line[:97] + "..."
                bullets.append(line)

        return bullets

    def _add_thankyou_slide(self, content: Dict, project_data: Dict):
        """Add thank you slide"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Background
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.PRIMARY_COLOR
        bg_shape.line.fill.background()

        # Thank you text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(12.333), Inches(2)
        )
        title_frame = title_box.text_frame

        p = title_frame.paragraphs[0]
        p.text = "Thank You!"
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER

        # Questions text
        q_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.5),
            Inches(12.333), Inches(1)
        )
        q_frame = q_box.text_frame

        p = q_frame.paragraphs[0]
        p.text = "Questions?"
        p.font.size = Pt(28)
        p.font.color.rgb = self.LIGHT_BG
        p.alignment = PP_ALIGN.CENTER

        # Contact info
        if content.get("contact"):
            contact_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(5.5),
                Inches(12.333), Inches(0.5)
            )
            contact_frame = contact_box.text_frame

            p = contact_frame.paragraphs[0]
            p.text = content.get("contact", "")
            p.font.size = Pt(16)
            p.font.color.rgb = self.LIGHT_BG
            p.alignment = PP_ALIGN.CENTER

    def add_architecture_slide(self, slide, architecture: Dict):
        """Add architecture diagram slide using shapes"""
        layers = architecture.get("layers", [])

        # Create layered boxes
        y_start = Inches(1.8)
        box_height = Inches(0.8)
        box_width = Inches(8)
        x_center = (self.SLIDE_WIDTH - box_width) / 2
        spacing = Inches(0.3)

        colors = [
            self.ACCENT_COLOR,
            self.SECONDARY_COLOR,
            self.PRIMARY_COLOR,
            RGBColor(100, 100, 100)
        ]

        for i, layer in enumerate(layers[:4]):
            y_pos = y_start + i * (box_height + spacing)
            color = colors[i % len(colors)]

            # Layer box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x_center, y_pos,
                box_width, box_height
            )
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.color.rgb = RGBColor(200, 200, 200)

            # Layer text
            text_frame = box.text_frame
            text_frame.paragraphs[0].text = layer.get("name", f"Layer {i+1}")
            text_frame.paragraphs[0].font.color.rgb = self.WHITE
            text_frame.paragraphs[0].font.size = Pt(16)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.word_wrap = True

            # Technology text
            if layer.get("technology"):
                p = text_frame.add_paragraph()
                p.text = layer.get("technology", "")
                p.font.size = Pt(12)
                p.font.color.rgb = self.LIGHT_BG
                p.alignment = PP_ALIGN.CENTER

            # Arrow between layers
            if i < len(layers) - 1:
                arrow_y = y_pos + box_height + Inches(0.05)
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.DOWN_ARROW,
                    (self.SLIDE_WIDTH / 2) - Inches(0.15),
                    arrow_y,
                    Inches(0.3), Inches(0.2)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.TEXT_COLOR
                arrow.line.fill.background()

    def add_code_slide(self, slide, code: str, language: str = ""):
        """Add code snippet slide"""
        # Code box with dark background
        code_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.5),
            Inches(12.333), Inches(5.5)
        )
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = RGBColor(30, 30, 30)
        code_box.line.color.rgb = RGBColor(60, 60, 60)

        # Code text
        text_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.7),
            Inches(11.933), Inches(5.1)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = False

        # Truncate code if too long
        code_lines = code.split('\n')[:20]

        for line in code_lines:
            p = text_frame.add_paragraph()
            p.text = line[:80]  # Truncate long lines
            p.font.name = 'Consolas'
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(200, 200, 200)


# Create instance
ppt_generator_v2 = PPTGeneratorV2()
