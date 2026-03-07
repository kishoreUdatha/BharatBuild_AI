from typing import Optional, Dict, Any
from pathlib import Path
import io
import zipfile
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER
from datetime import datetime

from app.core.config import settings
from app.core.logging_config import logger


class DocumentGenerator:
    """Generate various document formats (DOCX, PPTX, PDF)"""

    def __init__(self):
        self.temp_dir = settings.TEMP_DIR

    def generate_srs_docx(
        self,
        content: str,
        project_title: str,
        output_path: Optional[str] = None
    ) -> str:
        """
        Generate SRS document in DOCX format

        Args:
            content: SRS content
            project_title: Project title
            output_path: Output file path

        Returns:
            Path to generated document
        """
        if output_path is None:
            output_path = self.temp_dir / f"SRS_{project_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"

        doc = Document()

        # Title
        title = doc.add_heading(f'Software Requirements Specification', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Project Title
        subtitle = doc.add_heading(project_title, 1)
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Date
        date_para = doc.add_paragraph(f"Date: {datetime.now().strftime('%B %d, %Y')}")
        date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_page_break()

        # Parse and add content
        sections = content.split('\n\n')
        for section in sections:
            if section.strip():
                if section.startswith('#'):
                    # Heading
                    level = section.count('#', 0, 3)
                    text = section.lstrip('#').strip()
                    doc.add_heading(text, level)
                else:
                    # Paragraph
                    doc.add_paragraph(section.strip())

        doc.save(str(output_path))
        logger.info(f"Generated SRS DOCX: {output_path}")
        return str(output_path)

    def generate_report_docx(
        self,
        content: str,
        project_title: str,
        output_path: Optional[str] = None
    ) -> str:
        """Generate project report in DOCX format"""
        if output_path is None:
            output_path = self.temp_dir / f"Report_{project_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"

        doc = Document()

        # Title Page
        title = doc.add_heading('PROJECT REPORT', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_paragraph()
        project_title_para = doc.add_heading(project_title, 1)
        project_title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_paragraph()
        doc.add_paragraph()

        submitted_by = doc.add_paragraph('Submitted By:')
        submitted_by.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_paragraph()
        date_para = doc.add_paragraph(f"{datetime.now().strftime('%B %Y')}")
        date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_page_break()

        # Table of Contents placeholder
        doc.add_heading('Table of Contents', 1)
        doc.add_paragraph('(Auto-generated table of contents)')
        doc.add_page_break()

        # Content
        sections = content.split('\n\n')
        for section in sections:
            if section.strip():
                if section.startswith('#'):
                    level = min(section.count('#', 0, 3), 2)
                    text = section.lstrip('#').strip()
                    doc.add_heading(text, level)
                else:
                    doc.add_paragraph(section.strip())

        doc.save(str(output_path))
        logger.info(f"Generated Report DOCX: {output_path}")
        return str(output_path)

    def generate_ppt(
        self,
        slides_content: list,
        project_title: str,
        output_path: Optional[str] = None
    ) -> str:
        """
        Generate PowerPoint presentation

        Args:
            slides_content: List of dicts with 'title' and 'content'
            project_title: Project title
            output_path: Output file path

        Returns:
            Path to generated presentation
        """
        if output_path is None:
            output_path = self.temp_dir / f"Presentation_{project_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

        prs = Presentation()
        prs.slide_width = PptInches(10)
        prs.slide_height = PptInches(7.5)

        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = project_title
        subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"

        # Content Slides
        for slide_data in slides_content:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)

            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            title_shape.text = slide_data.get('title', 'Untitled')

            tf = body_shape.text_frame
            content = slide_data.get('content', '')

            if isinstance(content, list):
                for item in content:
                    p = tf.add_paragraph()
                    p.text = str(item)
                    p.level = 0
            else:
                tf.text = str(content)

        prs.save(str(output_path))
        logger.info(f"Generated PPT: {output_path}")
        return str(output_path)

    def generate_pdf(
        self,
        content: str,
        project_title: str,
        output_path: Optional[str] = None
    ) -> str:
        """Generate PDF document"""
        if output_path is None:
            output_path = self.temp_dir / f"Document_{project_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"

        doc = SimpleDocTemplate(str(output_path), pagesize=A4)
        story = []
        styles = getSampleStyleSheet()

        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=RGBColor(0, 0, 139),
            spaceAfter=30,
            alignment=TA_CENTER
        )

        story.append(Paragraph(project_title, title_style))
        story.append(Spacer(1, 0.2*inch))

        # Content
        sections = content.split('\n\n')
        for section in sections:
            if section.strip():
                if section.startswith('#'):
                    text = section.lstrip('#').strip()
                    story.append(Paragraph(text, styles['Heading2']))
                else:
                    story.append(Paragraph(section.strip(), styles['BodyText']))
                story.append(Spacer(1, 0.1*inch))

        doc.build(story)
        logger.info(f"Generated PDF: {output_path}")
        return str(output_path)

    def create_project_zip(
        self,
        files: Dict[str, str],
        project_title: str,
        output_path: Optional[str] = None
    ) -> str:
        """
        Create ZIP archive of project files

        Args:
            files: Dict of {filename: filepath}
            project_title: Project title
            output_path: Output ZIP path

        Returns:
            Path to generated ZIP file
        """
        if output_path is None:
            output_path = self.temp_dir / f"Project_{project_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"

        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for filename, filepath in files.items():
                if Path(filepath).exists():
                    zipf.write(filepath, arcname=filename)
                    logger.info(f"Added to ZIP: {filename}")

        logger.info(f"Created ZIP archive: {output_path}")
        return str(output_path)

    def generate_viva_qa_docx(
        self,
        qa_content: list,
        project_title: str,
        output_path: Optional[str] = None
    ) -> str:
        """
        Generate Viva Q&A document

        Args:
            qa_content: List of dicts with 'question' and 'answer'
            project_title: Project title
            output_path: Output file path

        Returns:
            Path to generated document
        """
        if output_path is None:
            output_path = self.temp_dir / f"Viva_QA_{project_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"

        doc = Document()

        # Title
        title = doc.add_heading('Viva Voce Preparation', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        subtitle = doc.add_heading(project_title, 1)
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_page_break()

        # Q&A
        for i, qa in enumerate(qa_content, 1):
            # Question
            q_para = doc.add_paragraph()
            q_run = q_para.add_run(f"Q{i}: {qa.get('question', '')}")
            q_run.bold = True
            q_run.font.size = Pt(12)

            # Answer
            a_para = doc.add_paragraph()
            a_run = a_para.add_run(f"A: {qa.get('answer', '')}")
            a_run.font.size = Pt(11)

            doc.add_paragraph()  # Spacing

        doc.save(str(output_path))
        logger.info(f"Generated Viva Q&A DOCX: {output_path}")
        return str(output_path)


    # ==================== NAAC/NBA ACCREDITATION METHODS ====================

    def generate_accreditation_docx(
        self,
        content: Dict[str, Any],
        title: str,
        output_path: Optional[str] = None
    ) -> bytes:
        """
        Generate NAAC/NBA accreditation document in DOCX format

        Args:
            content: Dictionary containing document content
            title: Document title

        Returns:
            Bytes of the generated Word document
        """
        doc = Document()

        # Title
        title_para = doc.add_heading(title, 0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Generation timestamp
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M")
        timestamp_para = doc.add_paragraph(f"Generated on: {timestamp}")
        timestamp_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_paragraph()

        # Header info
        doc.add_paragraph("NAAC Accreditation Document", style='Intense Quote')
        doc.add_paragraph("Generated by BharatBuild AI")
        doc.add_paragraph()

        # Process content recursively
        self._add_dict_content_to_word(doc, content, level=1)

        # Footer
        doc.add_paragraph()
        footer = doc.add_paragraph("---")
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("This document was generated by BharatBuild AI NAAC/NBA Accreditation System.")
        doc.add_paragraph("Please verify and customize the content as per your institution's requirements.")

        # Save to bytes
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        return file_stream.getvalue()

    def _add_dict_content_to_word(self, doc: Document, content: Any, level: int = 1):
        """Recursively add dictionary content to Word document."""
        if isinstance(content, dict):
            for key, value in content.items():
                if key in ['raw_response', 'parsed', 'metadata', 'error', 'success']:
                    continue  # Skip technical fields

                # Format key as heading
                heading_text = self._format_dict_key(key)
                if level <= 3:
                    doc.add_heading(heading_text, level)
                else:
                    para = doc.add_paragraph()
                    run = para.add_run(heading_text)
                    run.bold = True

                self._add_dict_content_to_word(doc, value, level + 1)

        elif isinstance(content, list):
            for i, item in enumerate(content):
                if isinstance(item, dict):
                    self._add_dict_content_to_word(doc, item, level)
                elif isinstance(item, str):
                    doc.add_paragraph(f"• {item}", style='List Bullet')
                else:
                    doc.add_paragraph(f"• {str(item)}", style='List Bullet')

        elif isinstance(content, str):
            if len(content) > 500:
                # Long text - add as multiple paragraphs
                paragraphs = content.split('\n')
                for para in paragraphs:
                    if para.strip():
                        doc.add_paragraph(para.strip())
            else:
                doc.add_paragraph(content)
        else:
            doc.add_paragraph(str(content))

    def _format_dict_key(self, key: str) -> str:
        """Format dictionary key as readable heading."""
        formatted = key.replace('_', ' ').replace('-', ' ')
        return formatted.title()

    def generate_accreditation_pdf(
        self,
        content: Dict[str, Any],
        title: str,
        output_path: Optional[str] = None
    ) -> bytes:
        """
        Generate NAAC/NBA accreditation document in PDF format

        Args:
            content: Dictionary containing document content
            title: Document title

        Returns:
            Bytes of the generated PDF document
        """
        buffer = io.BytesIO()
        doc_pdf = SimpleDocTemplate(
            buffer,
            pagesize=A4,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72
        )

        styles = getSampleStyleSheet()

        # Custom styles
        custom_title = ParagraphStyle(
            'AccredTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            alignment=TA_CENTER
        )

        custom_subtitle = ParagraphStyle(
            'AccredSubtitle',
            parent=styles['Normal'],
            fontSize=10,
            spaceAfter=20,
            alignment=TA_CENTER
        )

        story = []

        # Title
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M")
        story.append(Paragraph(title, custom_title))
        story.append(Paragraph(f"Generated on: {timestamp}", custom_subtitle))
        story.append(Spacer(1, 20))

        # Header
        story.append(Paragraph("NAAC Accreditation Document", styles['Heading3']))
        story.append(Paragraph("Generated by BharatBuild AI", styles['Normal']))
        story.append(Spacer(1, 20))

        # Process content
        self._add_dict_content_to_pdf(story, content, styles, level=1)

        # Footer
        story.append(Spacer(1, 30))
        story.append(Paragraph("---", custom_subtitle))
        story.append(Paragraph(
            "This document was generated by BharatBuild AI NAAC/NBA Accreditation System.",
            custom_subtitle
        ))

        # Build PDF
        doc_pdf.build(story)
        buffer.seek(0)
        return buffer.getvalue()

    def _add_dict_content_to_pdf(self, story: list, content: Any, styles, level: int = 1):
        """Recursively add dictionary content to PDF story."""
        if isinstance(content, dict):
            for key, value in content.items():
                if key in ['raw_response', 'parsed', 'metadata', 'error', 'success']:
                    continue

                heading_text = self._format_dict_key(key)

                if level == 1:
                    story.append(Paragraph(heading_text, styles['Heading1']))
                elif level == 2:
                    story.append(Paragraph(heading_text, styles['Heading2']))
                elif level == 3:
                    story.append(Paragraph(heading_text, styles['Heading3']))
                else:
                    story.append(Paragraph(f"<b>{heading_text}</b>", styles['Normal']))

                story.append(Spacer(1, 10))
                self._add_dict_content_to_pdf(story, value, styles, level + 1)

        elif isinstance(content, list):
            for item in content:
                if isinstance(item, dict):
                    self._add_dict_content_to_pdf(story, item, styles, level)
                elif isinstance(item, str):
                    story.append(Paragraph(f"• {item}", styles['Normal']))
                else:
                    story.append(Paragraph(f"• {str(item)}", styles['Normal']))
            story.append(Spacer(1, 10))

        elif isinstance(content, str):
            if len(content) > 500:
                paragraphs = content.split('\n')
                for para in paragraphs:
                    if para.strip():
                        story.append(Paragraph(para.strip(), styles['Normal']))
                        story.append(Spacer(1, 5))
            else:
                story.append(Paragraph(content, styles['Normal']))
                story.append(Spacer(1, 10))
        else:
            story.append(Paragraph(str(content), styles['Normal']))
            story.append(Spacer(1, 5))


# Create singleton instance
document_generator = DocumentGenerator()
