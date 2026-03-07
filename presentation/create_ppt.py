"""
BharatBuild.ai - College Demo Presentation Generator
Creates a professional PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand Colors
PRIMARY_ORANGE = RGBColor(255, 107, 0)  # #FF6B00
DARK_BG = RGBColor(15, 23, 42)  # #0F172A
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(148, 163, 184)  # #94A3B8
ACCENT_BLUE = RGBColor(59, 130, 246)  # #3B82F6
SUCCESS_GREEN = RGBColor(34, 197, 94)  # #22C55E

def set_slide_background(slide, color):
    """Set solid background color for slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_background(slide, DARK_BG)

    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "BharatBuild.ai"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "India's #1 AI-Powered Code Generation Platform"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Tagline
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.6))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Generate Complete Projects in Minutes, Not Weeks"
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER
    p.font.italic = True

def add_problem_slide(prs):
    """Slide 2: The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "The Problem Students Face"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

    problems = [
        ("2-3 Weeks", "Time wasted on final year projects"),
        ("No Documentation", "Existing tools don't generate SRS, Reports, PPTs"),
        ("Complex Setup", "Setting up development environment is hard"),
        ("Viva Struggles", "Students unprepared for project viva"),
        ("Outdated Tech", "Learning outdated frameworks"),
        ("Team Coordination", "Difficulty managing team contributions"),
    ]

    y_pos = 0.9
    for title, desc in problems:
        # Problem box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.6))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 41, 59)  # Slate-800
        box.line.fill.background()

        # Title
        text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.1), Inches(2.5), Inches(0.45))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(248, 113, 113)  # Red-400

        # Description
        desc_box = slide.shapes.add_textbox(Inches(3.5), Inches(y_pos + 0.12), Inches(5.5), Inches(0.45))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = LIGHT_GRAY

        y_pos += 0.72

def add_solution_slide(prs):
    """Slide 3: Our Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Our Solution: BharatBuild.ai"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

    solutions = [
        ("Complete in Minutes", "Generate full projects with AI in under 2 minutes"),
        ("Full Documentation", "SRS + UML + Report + PPT + Viva Q&A included"),
        ("Team of 3 Students", "One project, shared by 3 team members"),
        ("50+ Viva Questions", "Complete viva preparation with answers"),
        ("50+ Technologies", "React, Django, FastAPI, Flutter & more"),
        ("Production Ready", "Deploy-ready code with best practices"),
    ]

    y_pos = 0.9
    for title, desc in solutions:
        # Solution box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.6))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 41, 59)
        box.line.fill.background()

        # Checkmark + Title
        text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.1), Inches(3), Inches(0.45))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"+ {title}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = SUCCESS_GREEN

        # Description
        desc_box = slide.shapes.add_textbox(Inches(4), Inches(y_pos + 0.12), Inches(5), Inches(0.45))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE

        y_pos += 0.72

def add_why_not_free_slide(prs):
    """Slide: Why Not Free Tools?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Why Not Just Use Free Tools?"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Comparison headers
    headers = ["Task", "Free Tools (ChatGPT/Google)", "BharatBuild"]
    x_positions = [0.3, 2.8, 7.0]
    widths = [2.4, 4.0, 2.5]

    for header, x, w in zip(headers, x_positions, widths):
        cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(0.7), Inches(w), Inches(0.4))
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(51, 65, 85)
        cell.line.fill.background()

        text_box = slide.shapes.add_textbox(Inches(x), Inches(0.75), Inches(w), Inches(0.35))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    rows = [
        ["Complete Code", "Snippets only, manual assembly", "Full working project"],
        ["IEEE SRS Document", "Generic text, wrong format", "Perfect IEEE 830 format"],
        ["UML Diagrams", "Text description only", "Actual diagrams included"],
        ["Project Report", "5-10 pages max", "60-80 pages formatted"],
        ["Viva Questions", "Generic questions", "50+ project-specific Q&A"],
        ["Time Required", "2-3 weeks of work", "Under 5 minutes"],
        ["Integration", "Manual copy-paste", "Everything connected"],
    ]

    y_pos = 1.15
    for row in rows:
        for i, (cell_text, x, w) in enumerate(zip(row, x_positions, widths)):
            cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y_pos), Inches(w), Inches(0.5))
            cell.fill.solid()
            if i == 2:  # BharatBuild column
                cell.fill.fore_color.rgb = RGBColor(22, 101, 52)  # Green
            elif i == 1:  # Free tools column
                cell.fill.fore_color.rgb = RGBColor(127, 29, 29)  # Red
            else:
                cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
            cell.line.fill.background()

            text_box = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y_pos + 0.1), Inches(w - 0.1), Inches(0.35))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = cell_text
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

        y_pos += 0.55

    # Bottom message
    msg_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.05), Inches(9), Inches(0.4))
    tf = msg_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Free tools give you PIECES. BharatBuild gives you the COMPLETE PROJECT."
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

def add_time_comparison_slide(prs):
    """Slide: Time & Effort Comparison"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "The Real Cost of 'Free'"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Left side - Free Tools
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.9), Inches(4.5), Inches(4.2))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = RGBColor(127, 29, 29)  # Red
    left_card.line.fill.background()

    left_title = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(4.5), Inches(0.5))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Using Free Tools"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    free_steps = [
        "1. ChatGPT for code (3-4 days)",
        "2. Fix errors manually (2-3 days)",
        "3. Write SRS in Word (2 days)",
        "4. Draw UML in Lucidchart (1 day)",
        "5. Write Report (3-4 days)",
        "6. Make PPT in Canva (1 day)",
        "7. Search viva questions (1 day)",
        "8. Integrate everything (2 days)",
    ]

    y_pos = 1.5
    for step in free_steps:
        step_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos), Inches(4.1), Inches(0.35))
        tf = step_box.text_frame
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        y_pos += 0.38

    total_left = slide.shapes.add_textbox(Inches(0.4), Inches(4.6), Inches(4.5), Inches(0.4))
    tf = total_left.text_frame
    p = tf.paragraphs[0]
    p.text = "Total: 15-20 DAYS"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(248, 113, 113)
    p.alignment = PP_ALIGN.CENTER

    # Right side - BharatBuild
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(0.9), Inches(4.5), Inches(4.2))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = RGBColor(22, 101, 52)  # Green
    right_card.line.fill.background()

    right_title = slide.shapes.add_textbox(Inches(5.1), Inches(1.0), Inches(4.5), Inches(0.5))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Using BharatBuild"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    bb_steps = [
        "1. Describe your project",
        "2. AI generates everything:",
        "   - Complete source code",
        "   - IEEE SRS document",
        "   - UML diagrams",
        "   - 60-80 page report",
        "   - PPT presentation",
        "   - 50+ viva Q&A",
    ]

    y_pos = 1.5
    for step in bb_steps:
        step_box = slide.shapes.add_textbox(Inches(5.3), Inches(y_pos), Inches(4.1), Inches(0.35))
        tf = step_box.text_frame
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        y_pos += 0.38

    total_right = slide.shapes.add_textbox(Inches(5.1), Inches(4.6), Inches(4.5), Inches(0.4))
    tf = total_right.text_frame
    p = tf.paragraphs[0]
    p.text = "Total: 5 MINUTES"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SUCCESS_GREEN
    p.alignment = PP_ALIGN.CENTER

def add_features_slide(prs):
    """Slide 4: Key Features"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Powerful Features"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Left column features
    left_features = [
        "Full-Stack Code Generation",
        "50+ Tech Stacks Supported",
        "Live Code Preview",
        "Real-time AI Debugging",
        "Docker Containerization",
    ]

    # Right column features
    right_features = [
        "IEEE-Format SRS Documents",
        "UML Diagrams (Auto-generated)",
        "60-80 Page Project Reports",
        "Professional PPT Slides",
        "50+ Viva Questions & Answers",
    ]

    y_pos = 1.4
    for i, (left, right) in enumerate(zip(left_features, right_features)):
        # Left feature
        left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(4.3), Inches(0.65))
        left_box.fill.solid()
        left_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
        left_box.line.fill.background()

        text_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.15), Inches(4), Inches(0.5))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"* {left}"
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE

        # Right feature
        right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(y_pos), Inches(4.3), Inches(0.65))
        right_box.fill.solid()
        right_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
        right_box.line.fill.background()

        text_box = slide.shapes.add_textbox(Inches(5.4), Inches(y_pos + 0.15), Inches(4), Inches(0.5))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"* {right}"
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE

        y_pos += 0.8

def add_what_you_get_slide(prs):
    """Slide 5: What Students Get"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "What You Get For Rs.4499"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Complete Project Package for Team of 3 Students"
    p.font.size = Pt(20)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER

    items = [
        ("Complete Source Code", "Full-stack application with frontend + backend"),
        ("SRS Document", "IEEE 830 standard format (20+ pages)"),
        ("UML Diagrams", "Use Case, Class, Sequence, ER diagrams"),
        ("Project Report", "60-80 pages comprehensive documentation"),
        ("PPT Presentation", "Professional slides for project presentation"),
        ("Viva Q&A", "50+ questions with detailed answers"),
    ]

    y_pos = 1.5
    for title, desc in items:
        # Item box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y_pos), Inches(8.4), Inches(0.6))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 41, 59)
        box.line.fill.background()

        # Title
        text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.1), Inches(3), Inches(0.45))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"+ {title}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = SUCCESS_GREEN

        # Description
        desc_box = slide.shapes.add_textbox(Inches(4), Inches(y_pos + 0.12), Inches(5), Inches(0.45))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE

        y_pos += 0.7

def add_tech_stack_slide(prs):
    """Slide 6: Technology Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "50+ Supported Technologies"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

    categories = [
        ("Frontend", "React, Next.js, Vue, Angular, Svelte, Tailwind CSS"),
        ("Backend", "FastAPI, Django, Flask, Express, Spring Boot, Go, Rust"),
        ("Mobile", "React Native, Flutter, Kotlin, Swift"),
        ("Database", "PostgreSQL, MongoDB, MySQL, Redis, Firebase"),
        ("DevOps", "Docker, AWS, Vercel, GitHub Actions, Nginx"),
    ]

    y_pos = 1.3
    for category, techs in categories:
        # Category box
        cat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y_pos), Inches(2), Inches(0.7))
        cat_box.fill.solid()
        cat_box.fill.fore_color.rgb = PRIMARY_ORANGE
        cat_box.line.fill.background()

        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.15), Inches(2), Inches(0.5))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = category
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Tech list
        tech_box = slide.shapes.add_textbox(Inches(3), Inches(y_pos + 0.15), Inches(6), Inches(0.5))
        tf = tech_box.text_frame
        p = tf.paragraphs[0]
        p.text = techs
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE

        y_pos += 0.85

def add_agent_architecture_slide(prs):
    """Slide 7: AI Agent Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "25+ AI Agents Working Together"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Workflow
    workflow_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.6))
    tf = workflow_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Planner -> Writer -> Runner -> Fixer -> Documenter"
    p.font.size = Pt(20)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER

    agents = [
        ("Core Agents", ["PlannerAgent - Understands requirements",
                         "WriterAgent - Generates code",
                         "FixerAgent - Debugs errors",
                         "RunnerAgent - Executes code"]),
        ("Academic Agents", ["SRSAgent - IEEE 830 documents",
                             "UMLAgent - Diagrams",
                             "ReportAgent - Project reports",
                             "PPTAgent - Presentations",
                             "VivaAgent - Q&A preparation"]),
    ]

    x_pos = 0.5
    for category, agent_list in agents:
        # Category header
        header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.8), Inches(4.5), Inches(0.5))
        header.fill.solid()
        header.fill.fore_color.rgb = PRIMARY_ORANGE
        header.line.fill.background()

        text_box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.9), Inches(4.5), Inches(0.4))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = category
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Agent list
        y_pos = 2.5
        for agent in agent_list:
            agent_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos + 0.1), Inches(y_pos), Inches(4.3), Inches(0.5))
            agent_box.fill.solid()
            agent_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
            agent_box.line.fill.background()

            text_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(y_pos + 0.1), Inches(4.1), Inches(0.4))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"* {agent}"
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE

            y_pos += 0.55

        x_pos += 4.8

def add_pricing_slide(prs):
    """Slide 8: Pricing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Simple, Affordable Pricing"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Main pricing card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(1.2), Inches(6), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(30, 41, 59)
    card.line.fill.background()

    # Price
    price_box = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(6), Inches(1))
    tf = price_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Rs. 4,499"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Per project
    per_box = slide.shapes.add_textbox(Inches(2), Inches(2.4), Inches(6), Inches(0.5))
    tf = per_box.text_frame
    p = tf.paragraphs[0]
    p.text = "per complete project"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Team info
    team_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(0.5))
    tf = team_box.text_frame
    p = tf.paragraphs[0]
    p.text = "For a Team of 3 Students"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Per student calculation
    calc_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(0.5))
    tf = calc_box.text_frame
    p = tf.paragraphs[0]
    p.text = "= Rs. 1,500 per student"
    p.font.size = Pt(20)
    p.font.color.rgb = SUCCESS_GREEN
    p.alignment = PP_ALIGN.CENTER

    # What's included
    inc_box = slide.shapes.add_textbox(Inches(2), Inches(4.1), Inches(6), Inches(0.5))
    tf = inc_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Code + SRS + UML + Report + PPT + Viva Q&A"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

def add_stats_slide(prs):
    """Slide 9: Traction & Stats"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Our Traction"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

    stats = [
        ("10,000+", "Projects Built"),
        ("5,000+", "Active Students"),
        ("3,000+", "Developers"),
        ("500+", "Startups"),
        ("50+", "Institutions"),
        ("4.8/5", "User Rating"),
    ]

    # Create 2x3 grid
    positions = [
        (0.5, 1.5), (3.5, 1.5), (6.5, 1.5),
        (0.5, 3.5), (3.5, 3.5), (6.5, 3.5),
    ]

    for (x, y), (number, label) in zip(positions, stats):
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.8), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(30, 41, 59)
        card.line.fill.background()

        # Number
        num_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.3), Inches(2.8), Inches(0.7))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_ORANGE
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.95), Inches(2.8), Inches(0.4))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

def add_college_benefits_slide(prs):
    """Slide 10: Benefits for Colleges"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Benefits for Your College"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

    benefits = [
        ("For Students", [
            "Save 2-3 weeks on final year projects",
            "Complete documentation included",
            "Learn industry-standard practices",
            "Build real portfolio projects",
        ]),
        ("For Faculty", [
            "Better quality student projects",
            "Standardized documentation format",
            "Students learn modern tech stacks",
            "Focus on teaching, not struggles",
        ]),
        ("For Institution", [
            "Higher project completion rates",
            "Better placement statistics",
            "Modern AI-powered learning",
            "Competitive advantage",
        ]),
    ]

    x_pos = 0.4
    for title, items in benefits:
        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.2), Inches(3), Inches(0.5))
        header.fill.solid()
        header.fill.fore_color.rgb = ACCENT_BLUE
        header.line.fill.background()

        text_box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.3), Inches(3), Inches(0.4))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Items
        y_pos = 1.9
        for item in items:
            item_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos + 0.05), Inches(y_pos), Inches(2.9), Inches(0.55))
            item_box.fill.solid()
            item_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
            item_box.line.fill.background()

            text_box = slide.shapes.add_textbox(Inches(x_pos + 0.15), Inches(y_pos + 0.1), Inches(2.7), Inches(0.45))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"+ {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE

            y_pos += 0.65

        x_pos += 3.2

def add_how_it_works_slide(prs):
    """Slide 11: How It Works"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "How It Works"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

    steps = [
        ("1", "Describe Your Project", "Tell AI what you want to build"),
        ("2", "AI Generates Everything", "Code, docs, diagrams in minutes"),
        ("3", "Review & Customize", "Make changes if needed"),
        ("4", "Download & Submit", "Get all files ready for submission"),
    ]

    y_pos = 1.3
    for num, title, desc in steps:
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(y_pos), Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY_ORANGE
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.1), Inches(0.7), Inches(0.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Title
        title_box = slide.shapes.add_textbox(Inches(2), Inches(y_pos + 0.05), Inches(3), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Description
        desc_box = slide.shapes.add_textbox(Inches(2), Inches(y_pos + 0.4), Inches(6), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_GRAY

        y_pos += 1.0

def add_demo_slide(prs):
    """Slide 12: Live Demo"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Live Demo"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Watch AI Build a Complete Project in Minutes"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # URL
    url_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.6))
    tf = url_box.text_frame
    p = tf.paragraphs[0]
    p.text = "bharatbuild.ai"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER

def add_cta_slide(prs):
    """Slide 13: Call to Action"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Start Building Today!"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Pricing reminder
    price_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(0.8))
    tf = price_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Rs. 4,499 per project (Team of 3)"
    p.font.size = Pt(28)
    p.font.color.rgb = SUCCESS_GREEN
    p.alignment = PP_ALIGN.CENTER

    # What you get
    get_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(0.5))
    tf = get_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Code + SRS + UML + Report + PPT + 50+ Viva Questions"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # URL Box
    url_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(3.8), Inches(5), Inches(0.9))
    url_card.fill.solid()
    url_card.fill.fore_color.rgb = PRIMARY_ORANGE
    url_card.line.fill.background()

    url_box = slide.shapes.add_textbox(Inches(2.5), Inches(4), Inches(5), Inches(0.7))
    tf = url_box.text_frame
    p = tf.paragraphs[0]
    p.text = "bharatbuild.ai"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_thank_you_slide(prs):
    """Slide 14: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    # Thank you
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Questions
    q_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.6))
    tf = q_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1.2))
    tf = contact_box.text_frame

    p = tf.paragraphs[0]
    p.text = "Website: bharatbuild.ai"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Email: contact@bharatbuild.ai"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Location: Hyderabad, India"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

def main():
    """Generate the presentation"""
    # Create presentation (16:9 aspect ratio)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Add all slides
    add_title_slide(prs)              # 1 - Title
    add_problem_slide(prs)            # 2 - Problems students face
    add_solution_slide(prs)           # 3 - Our solution
    add_why_not_free_slide(prs)       # 4 - WHY NOT FREE TOOLS (NEW - addresses objection)
    add_time_comparison_slide(prs)    # 5 - TIME COMPARISON (NEW - shows real cost)
    add_what_you_get_slide(prs)       # 6 - What you get for Rs.4499
    add_features_slide(prs)           # 7 - Features
    add_tech_stack_slide(prs)         # 8 - Tech stack
    add_agent_architecture_slide(prs) # 9 - AI agents
    add_pricing_slide(prs)            # 10 - Pricing
    add_college_benefits_slide(prs)   # 11 - Benefits for colleges
    add_how_it_works_slide(prs)       # 12 - How it works
    add_demo_slide(prs)               # 13 - Live demo
    add_cta_slide(prs)                # 14 - Call to action
    add_thank_you_slide(prs)          # 15 - Thank you

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "BharatBuild_College_Demo.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    main()
