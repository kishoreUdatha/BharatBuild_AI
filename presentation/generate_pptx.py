"""
BharatBuild AI Pitch Deck Generator
Generates a PowerPoint presentation following the ChatBucket format
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
PRIMARY_COLOR = RGBColor(0, 82, 204)  # Blue
SECONDARY_COLOR = RGBColor(255, 107, 0)  # Orange
DARK_COLOR = RGBColor(30, 30, 30)  # Dark gray
WHITE_COLOR = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(245, 247, 250)

def add_title_slide(prs, title, subtitle, tagline=None):
    """Slide 1: Title Slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_COLOR
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Tagline
    if tagline:
        tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12), Inches(0.6))
        tf = tag_box.text_frame
        p = tf.paragraphs[0]
        p.text = tagline
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(180, 180, 180)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, subtitle=None):
    """Generic content slide with bullets"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_COLOR
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR

    # Subtitle if provided
    y_start = Inches(1.5)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = SECONDARY_COLOR
        p.font.italic = True
        y_start = Inches(2.0)

    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(0.5), y_start, Inches(12), Inches(5))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_COLOR
        p.space_after = Pt(12)

    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    """Two column layout"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_COLOR
    title_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # Left bullets
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.5), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_COLOR
        p.space_after = Pt(8)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # Right bullets
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(5.5), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_COLOR
        p.space_after = Pt(8)

    return slide

def add_stats_slide(prs, title, stats):
    """Stats/metrics slide with boxes"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_COLOR
    title_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR

    # Stats boxes
    num_stats = len(stats)
    box_width = Inches(2.8)
    total_width = num_stats * 2.8 + (num_stats - 1) * 0.3
    start_x = (13.333 - total_width) / 2

    for i, (value, label) in enumerate(stats):
        x = Inches(start_x + i * 3.1)

        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.5), box_width, Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(2)

        # Value
        val_box = slide.shapes.add_textbox(x, Inches(2.8), box_width, Inches(1))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER

        # Label
        lab_box = slide.shapes.add_textbox(x, Inches(3.8), box_width, Inches(0.8))
        tf = lab_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_COLOR
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_table_slide(prs, title, headers, rows):
    """Table slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PRIMARY_COLOR
    title_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR

    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1

    table_width = Inches(12)
    table_height = Inches(0.5 * num_rows)
    left = Inches(0.65)
    top = Inches(1.8)

    table = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height).table

    # Set column widths
    col_width = Inches(12 / num_cols)
    for i in range(num_cols):
        table.columns[i].width = col_width

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE_COLOR
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_COLOR
            p.alignment = PP_ALIGN.CENTER
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG

    return slide

def add_quote_slide(prs, quote, author, role):
    """Quote/testimonial slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()

    # Quote mark
    quote_mark = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(1), Inches(1))
    tf = quote_mark.text_frame
    p = tf.paragraphs[0]
    p.text = '"'
    p.font.size = Pt(120)
    p.font.color.rgb = PRIMARY_COLOR

    # Quote text
    quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(2))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = DARK_COLOR

    # Author
    author_box = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(10), Inches(0.5))
    tf = author_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"— {author}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR

    # Role
    role_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10), Inches(0.5))
    tf = role_box.text_frame
    p = tf.paragraphs[0]
    p.text = role
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(100, 100, 100)

    return slide

def add_closing_slide(prs, company, tagline, contact_info):
    """Closing slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_COLOR
    bg.line.fill.background()

    # Company name
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = company
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Tagline
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(12), Inches(0.8))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = tagline
    p.font.size = Pt(24)
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12), Inches(1.5))
    tf = contact_box.text_frame
    for i, info in enumerate(contact_info):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = info
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(180, 180, 180)
        p.alignment = PP_ALIGN.CENTER

    return slide

# ============================================
# CREATE THE PITCH DECK
# ============================================

print("Creating BharatBuild AI Pitch Deck...")

# Slide 1: Title
add_title_slide(
    prs,
    "BharatBuild AI",
    "AI-Powered Academic Excellence & Accreditation for Bharat",
    "India's First AI Platform for Students, Colleges & Accreditation"
)

# Slide 2: Vision Statement
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = PRIMARY_COLOR
bg.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Democratizing AI for"
p.font.size = Pt(36)
p.font.color.rgb = WHITE_COLOR
p.alignment = PP_ALIGN.CENTER

title_box2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(12), Inches(1.5))
tf = title_box2.text_frame
p = tf.paragraphs[0]
p.text = "40 Million Indian Students"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE_COLOR
p.alignment = PP_ALIGN.CENTER

sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12), Inches(0.8))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "Building the Bolt.new + NAAC Automation for Bharat"
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(200, 220, 255)
p.alignment = PP_ALIGN.CENTER

# Slide 3: Problem Statement
add_content_slide(
    prs,
    "The Problem",
    [
        "4.5 Million engineering students graduate annually in India",
        "95% cannot access international AI tools (no intl. credit cards)",
        "Bolt.new/Cursor costs ₹1,700/month — unaffordable for Indian students",
        "1,100+ colleges face NAAC accreditation — 80% do it manually",
        "Zero platforms combine project generation + accreditation automation"
    ],
    "Indian students and colleges are locked out of modern AI development tools"
)

# Slide 4: Who is Affected (ICP)
add_content_slide(
    prs,
    "Who We Serve",
    [
        "B.Tech/B.E. Students — Final year projects, documentation, viva prep",
        "CS/IT Faculty — Project management, OBE compliance, batch tracking",
        "IQAC Coordinators — NAAC/NBA accreditation documentation",
        "College Management — Rankings, compliance, institutional quality"
    ],
    "Ideal Customer Profile"
)

# Slide 5: Solution Overview
add_content_slide(
    prs,
    "Our Solution",
    [
        "AI-powered full-stack project generation (code + documentation)",
        "Automated SRS, UML diagrams, project reports, PPT, viva Q&A",
        "Real-time code execution with live preview (like Bolt.new)",
        "NAAC/NBA accreditation automation (7 criteria, 700 marks)",
        "OBE compliance engine (CO-PO mapping, attainment calculation)",
        "UPI/Razorpay payments — no international credit card needed"
    ],
    "The Bolt.new + NAAC automation for India — at 85% lower cost"
)

# Slide 6: Key Features
add_two_column_slide(
    prs,
    "Platform Features",
    "For Students",
    [
        "One-click project generation",
        "50+ page project reports",
        "SRS, UML, ER diagrams",
        "PowerPoint presentations",
        "Viva Q&A preparation (50+ questions)",
        "Real-time code editor & preview"
    ],
    "For Colleges",
    [
        "NAAC 7-criteria automation",
        "OBE compliance dashboard",
        "Faculty batch management",
        "CO-PO mapping & attainment",
        "IQAC workflow approvals",
        "Plagiarism detection"
    ]
)

# Slide 7: Traction Stats
add_stats_slide(
    prs,
    "Traction & Validation",
    [
        ("500+", "Beta Students"),
        ("12", "Pilot Colleges"),
        ("78%", "Weekly Active"),
        ("3", "MOUs Signed")
    ]
)

# Slide 8: Testimonial
add_quote_slide(
    prs,
    "We reduced our NAAC preparation time from 6 months to 6 weeks. BharatBuild AI automated 80% of our documentation work.",
    "Dr. Ramesh Kumar",
    "IQAC Coordinator, ABC Engineering College"
)

# Slide 9: How It Works
add_content_slide(
    prs,
    "How It Works",
    [
        "Step 1: Choose Your Mode — Student | Developer | College | Founder",
        "Step 2: Describe Your Project — AI understands requirements",
        "Step 3: 20+ Agents Generate Everything — Code, docs, diagrams",
        "Step 4: Review & Customize — Edit in real-time code editor",
        "Step 5: Download & Deploy — Production-ready code + formatted docs"
    ]
)

# Slide 10: Market Opportunity
add_content_slide(
    prs,
    "Market Opportunity",
    [
        "TAM: ₹45,000 Cr — Indian EdTech market (2025)",
        "SAM: ₹4,500 Cr — Engineering education tools",
        "SOM: ₹150 Cr — Year 3 target (500 colleges + 2L students)",
        "Why Now: NEP 2020 mandates OBE, NAAC 2024 requires digital docs",
        "Growth: B2B2C model with network effects"
    ],
    "4.5M engineering students + 40,000+ colleges in India"
)

# Slide 11: TAM/SAM/SOM Visual
add_stats_slide(
    prs,
    "Addressable Market",
    [
        ("₹45,000 Cr", "TAM - EdTech"),
        ("₹4,500 Cr", "SAM - Eng. Tools"),
        ("₹150 Cr", "SOM - Year 3")
    ]
)

# Slide 12: Current Stage
add_content_slide(
    prs,
    "Current Stage",
    [
        "MVP Live (v1.2) — 5 core modules deployed",
        "150,000+ lines of production code",
        "20+ AI agents for specialized tasks",
        "Docker-based code execution sandbox",
        "Razorpay payment integration complete",
        "Bootstrapped: ₹35 Lakhs invested"
    ],
    "Product Stage: Beta with 12 colleges"
)

# Slide 13: Revenue Model - B2C
add_table_slide(
    prs,
    "Revenue Model — Students (B2C)",
    ["Plan", "Price/Month", "Features"],
    [
        ["Free", "₹0", "1 project, basic docs"],
        ["Pro", "₹99", "Unlimited projects, all docs"],
        ["Premium", "₹299", "Priority generation, viva prep"]
    ]
)

# Slide 14: Revenue Model - B2B
add_table_slide(
    prs,
    "Revenue Model — Colleges (B2B)",
    ["Plan", "Price/Student/Year", "Features"],
    [
        ["Starter", "₹49", "Basic project generation"],
        ["Professional", "₹99", "Full suite + faculty dashboard"],
        ["Enterprise", "₹149", "NAAC automation + API access"]
    ]
)

# Slide 15: Competitor Matrix
add_table_slide(
    prs,
    "Competitive Advantage",
    ["Feature", "BharatBuild", "Bolt.new", "Cursor", "Replit"],
    [
        ["Price (₹/month)", "₹99", "₹1,700", "₹1,700", "₹600"],
        ["Indian Payments", "✓", "✗", "✗", "✗"],
        ["Academic Docs", "✓", "✗", "✗", "✗"],
        ["NAAC/NBA Support", "✓", "✗", "✗", "✗"],
        ["Indian Servers", "✓", "✗", "✗", "✗"]
    ]
)

# Slide 16: Financial Projections
add_table_slide(
    prs,
    "Financial Projections",
    ["Metric", "Year 1", "Year 2", "Year 3"],
    [
        ["Revenue", "₹1.5 Cr", "₹8 Cr", "₹25 Cr"],
        ["Students", "25,000", "1,00,000", "2,00,000"],
        ["Colleges", "50", "200", "500"],
        ["Break-even", "—", "Month 18", "Profitable"]
    ]
)

# Slide 17: GTM Strategy
add_content_slide(
    prs,
    "Go-To-Market Strategy",
    [
        "Phase 1 (Year 1): Karnataka & Telangana — 50 colleges, 25K students",
        "Phase 2 (Year 2): South India — 200 colleges, 1L students",
        "Phase 3 (Year 3): Pan-India — 500 colleges, 2L+ students",
        "Channels: Direct sales + college referrals + channel partners",
        "Partnerships: University ERPs, book publishers, exam prep companies"
    ]
)

# Slide 18: The Ask
add_content_slide(
    prs,
    "The Ask",
    [
        "Raising: ₹5 Crore (Seed Round)",
        "Equity: 12-15%",
        "Valuation: ₹35 Crore (post-money)",
        "Runway: 18 months to Series A"
    ],
    "Investment Scope"
)

# Slide 19: Use of Funds
add_stats_slide(
    prs,
    "Use of Funds",
    [
        ("40%", "AI Infrastructure"),
        ("30%", "Team Expansion"),
        ("20%", "Go-to-Market"),
        ("10%", "Operations")
    ]
)

# Slide 20: What We're Looking For
add_content_slide(
    prs,
    "What We're Looking For",
    [
        "Strategic Partners: EdTech companies, University ERP providers",
        "Angel Investors: EdTech/SaaS experience, college networks",
        "VCs: India B2B SaaS focus, education sector expertise",
        "Value Beyond Capital: College intros, GTM expertise, regulatory guidance"
    ]
)

# Slide 21: Team
add_content_slide(
    prs,
    "The Team",
    [
        "Founder/CEO — 8+ years EdTech & SaaS, scaled to 500K users",
        "Co-founder/CTO — 10+ years AI/ML, built 5 production AI systems",
        "Co-founder/COO — 6+ years ops, managed 200+ college relationships",
        "",
        "Why We Win:",
        "• Deep domain expertise in Indian engineering education",
        "• Technical depth — built AI systems, not just API wrappers",
        "• Existing network of 50+ colleges",
        "• Execution speed — MVP in 4 months, 12 pilots in 6 months"
    ]
)

# Slide 22: Closing
add_closing_slide(
    prs,
    "BharatBuild AI",
    "Let's build the future of Indian engineering education together.",
    [
        "founders@bharatbuild.ai",
        "+91 XXXXX XXXXX",
        "www.bharatbuild.ai"
    ]
)

# Save the presentation
output_path = "BharatBuild_AI_Pitch_Deck.pptx"
prs.save(output_path)
print(f"Pitch deck saved to: {output_path}")
print("Done! 22 slides created.")
