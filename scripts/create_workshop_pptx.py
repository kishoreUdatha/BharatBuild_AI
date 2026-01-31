from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(15, 23, 42)
PURPLE = RGBColor(99, 102, 241)
LIGHT_PURPLE = RGBColor(165, 180, 252)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(148, 163, 184)
YELLOW = RGBColor(251, 191, 36)

def add_title_slide(title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = PURPLE
    background.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(title, bullets=None, two_cols=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = LIGHT_PURPLE

    if bullets:
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(24)
            p.font.color.rgb = WHITE
            p.space_after = Pt(12)

    if two_cols:
        left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.8), Inches(5.5))
        tf = left_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = two_cols[0][0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = YELLOW
        for item in two_cols[0][1]:
            p = tf.add_paragraph()
            p.text = "  " + item
            p.font.size = Pt(22)
            p.font.color.rgb = WHITE
            p.space_after = Pt(8)

        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5))
        tf = right_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = two_cols[1][0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = YELLOW
        for item in two_cols[1][1]:
            p = tf.add_paragraph()
            p.text = "  " + item
            p.font.size = Pt(22)
            p.font.color.rgb = WHITE
            p.space_after = Pt(8)

    return slide

def add_section_slide(title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(30, 58, 95)
    background.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    return slide

# ===== CREATE SLIDES =====

# Slide 1: Title
add_title_slide("BharatBuild AI", "Project Generation Workshop\n\n2-Day Hands-on Program\n10:00 AM - 4:00 PM")

# Slide 2: Workshop Overview
add_content_slide("Workshop Overview", two_cols=[
    ("Day 1: Fundamentals", [
        "Introduction to AI Generation",
        "Writing Effective Prompts",
        "Generate Simple Projects",
        "Modify with AI Chat"
    ]),
    ("Day 2: Advanced", [
        "Full-Stack Projects",
        "Complex Applications",
        "Download & Deploy",
        "Project Showcase"
    ])
])

# Slide 3: What is BharatBuild AI
add_content_slide("What is BharatBuild AI?", bullets=[
    "AI-powered platform that generates complete code from natural language",
    "",
    "Describe your project in English - Get full source code",
    "",
    "Live preview - See your app running instantly",
    "",
    "AI modifications - Ask AI to add features & fix bugs",
    "",
    "Download complete project files - Deploy anywhere"
])

# Slide 4: Supported Technologies
add_content_slide("Supported Technologies", two_cols=[
    ("Frontend", [
        "React / Next.js",
        "Vue.js / Nuxt.js",
        "Angular",
        "HTML / CSS / JavaScript",
        "Tailwind CSS"
    ]),
    ("Backend & Database", [
        "Node.js / Express",
        "Python / FastAPI / Django",
        "Java / Spring Boot",
        "PostgreSQL / MySQL",
        "MongoDB"
    ])
])

# Slide 5: Day 1 Section
add_section_slide("DAY 1", "Fundamentals & Simple Projects")

# Slide 6: Getting Started
add_content_slide("Getting Started", bullets=[
    "Step 1:  Open bharatbuild.ai in your browser",
    "",
    "Step 2:  Login with your credentials",
    "",
    "Step 3:  Click 'New Project'",
    "",
    "Step 4:  Enter your project description",
    "",
    "Step 5:  Click 'Generate' and watch the magic!"
])

# Slide 7: Bad vs Good Prompts
add_content_slide("Writing Effective Prompts", bullets=[
    "BAD PROMPT:",
    "    'Make me a website'",
    "",
    "GOOD PROMPT:",
    "    'Create a personal portfolio website with React and",
    "     Tailwind CSS. Include sections for About Me, Skills,",
    "     Projects with image cards, and Contact form.",
    "     Use dark theme with purple accents.'"
])

# Slide 8: Prompt Structure
add_content_slide("Prompt Structure", bullets=[
    "Project Type:   'Create a todo app' / 'Build an e-commerce site'",
    "",
    "Tech Stack:     'using React and Node.js' / 'with Next.js'",
    "",
    "Features:       'Include login, dashboard, and reports'",
    "",
    "Design:         'Use dark theme with blue accents'",
    "",
    "Pages:          'Include Home, About, Products, Contact pages'"
])

# Slide 9: Live Demo
add_section_slide("LIVE DEMO", "Generating a Portfolio Website")

# Slide 10: Hands-on Exercise 1
add_content_slide("Hands-on Exercise 1", bullets=[
    "Generate your first project! Choose one:",
    "",
    "    Todo App - Add, complete, delete tasks",
    "",
    "    Calculator - Basic arithmetic operations",
    "",
    "    Weather App - Show weather for a city",
    "",
    "    Landing Page - Simple business page",
    "",
    "Time: 30 minutes"
])

# Slide 11: Understanding Generated Code
add_content_slide("Understanding Generated Code", two_cols=[
    ("File Structure", [
        "src/ - Source code folder",
        "components/ - UI components",
        "pages/ - Route pages",
        "styles/ - CSS files",
        "package.json - Dependencies"
    ]),
    ("Key Actions", [
        "Browse files in editor",
        "Click to open any file",
        "Preview runs automatically",
        "Edit code directly",
        "Ask AI for changes"
    ])
])

# Slide 12: Modifying with AI Chat
add_content_slide("Modifying with AI Chat", bullets=[
    "After generation, ask AI to make changes:",
    "",
    "    'Add a dark/light theme toggle button'",
    "",
    "    'Change the primary color from blue to purple'",
    "",
    "    'Add a new Services page with 6 cards'",
    "",
    "    'Fix the mobile navigation menu'",
    "",
    "    'Add form validation with error messages'"
])

# Slide 13: Day 2 Section
add_section_slide("DAY 2", "Advanced Projects & Deployment")

# Slide 14: Full-Stack Projects
add_content_slide("Full-Stack Project Generation", bullets=[
    "Example: Full-Stack Blog Application",
    "",
    "    Frontend: Next.js with Tailwind CSS",
    "    Backend: Node.js with Express",
    "    Database: MongoDB",
    "",
    "    Features:",
    "    - User authentication (login/register)",
    "    - Create, edit, delete blog posts",
    "    - Categories, tags, and comments",
    "    - Admin dashboard"
])

# Slide 15: Complex Project Examples
add_content_slide("Complex Project Examples", two_cols=[
    ("Business Apps", [
        "E-Commerce Website",
        "Hospital Management",
        "Student Portal",
        "Inventory System"
    ]),
    ("More Ideas", [
        "Restaurant Ordering",
        "HR Management",
        "Analytics Dashboard",
        "Chat Application"
    ])
])

# Slide 16: Hands-on Exercise 2
add_content_slide("Hands-on Exercise 2 - Final Project", bullets=[
    "Build your final project! Choose one:",
    "",
    "    Student Management System",
    "",
    "    E-Commerce Website",
    "",
    "    Blog Platform with Admin Panel",
    "",
    "    Task Manager with Teams",
    "",
    "    Your Own Idea!",
    "",
    "Time: 60 minutes"
])

# Slide 17: Download & Deploy
add_content_slide("Download & Deploy", bullets=[
    "Step 1:  Click 'Download' to get ZIP file",
    "",
    "Step 2:  Extract and open in VS Code",
    "",
    "Step 3:  Run 'npm install' to install dependencies",
    "",
    "Step 4:  Run 'npm run dev' to start locally",
    "",
    "Step 5:  Deploy to Vercel / Netlify / Railway",
    "",
    "All platforms offer FREE hosting!"
])

# Slide 18: Tips & Best Practices
add_content_slide("Tips & Best Practices", bullets=[
    "  Be specific in your prompts",
    "",
    "  Mention the tech stack you want",
    "",
    "  List all features clearly",
    "",
    "  Describe the design/theme",
    "",
    "  Use AI chat for incremental changes",
    "",
    "  Review and understand the generated code",
    "",
    "  Don't use vague prompts like 'make it better'"
])

# Slide 19: Q&A
add_section_slide("Questions?", "Let's discuss!")

# Slide 20: Thank You
add_title_slide("Thank You!", "Happy Building with BharatBuild AI\n\nbharatbuild.ai")

# Save
output_path = r"C:\Users\Lekhana\Kishore_projects\BharatBuild_AI\docs\BharatBuild_Workshop.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
