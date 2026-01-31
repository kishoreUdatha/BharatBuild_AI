"""
BharatBuild.ai - Premium PowerPoint Presentation
Professional design with modern themes and compelling content
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# Create widescreen presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ========== PREMIUM COLOR PALETTE ==========
# Primary gradient colors
GRADIENT_START = RGBColor(6, 182, 212)      # Cyan 500
GRADIENT_END = RGBColor(59, 130, 246)       # Blue 500
ACCENT_PURPLE = RGBColor(139, 92, 246)      # Violet 500

# Background colors
BG_DARK = RGBColor(2, 6, 23)                # Very dark blue
BG_CARD = RGBColor(15, 23, 42)              # Slate 900
BG_CARD_LIGHT = RGBColor(30, 41, 59)        # Slate 800

# Text colors
TEXT_WHITE = RGBColor(255, 255, 255)
TEXT_GRAY = RGBColor(148, 163, 184)         # Slate 400
TEXT_LIGHT = RGBColor(226, 232, 240)        # Slate 200

# Accent colors
ACCENT_GREEN = RGBColor(34, 197, 94)        # Green 500
ACCENT_RED = RGBColor(239, 68, 68)          # Red 500
ACCENT_ORANGE = RGBColor(251, 146, 60)      # Orange 400
ACCENT_YELLOW = RGBColor(250, 204, 21)      # Yellow 400

def add_dark_background(slide):
    """Add premium dark background"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_accent_shape(slide, left, top, width, height, color):
    """Add colored accent shape"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height, size=18, color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT):
    """Add text box with styling"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_bullet_list(slide, items, left, top, width, height, size=16, color=TEXT_GRAY, bullet_color=ACCENT_GREEN):
    """Add bullet list"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_before = Pt(8)
        p.space_after = Pt(4)
    return txBox

def add_card(slide, left, top, width, height, border_color=None):
    """Add styled card"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(2)
    else:
        card.line.color.rgb = RGBColor(51, 65, 85)
        card.line.width = Pt(1)
    return card

def add_icon_circle(slide, left, top, size, color, text=""):
    """Add circular icon placeholder"""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    if text:
        add_text(slide, text, left, top + Inches(0.05), size, size,
                size=24, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    return circle

# =====================================================
# SLIDE 1: TITLE SLIDE - Hero Section (Clean Design)
# =====================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide1)

# Top accent bar
top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = GRADIENT_START
top_bar.line.fill.background()

# Logo/Brand area - centered box
brand_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.2), Inches(4.3), Inches(0.6))
brand_box.fill.solid()
brand_box.fill.fore_color.rgb = GRADIENT_START
brand_box.line.fill.background()
add_text(slide1, "AI-POWERED DEVELOPMENT", Inches(4.5), Inches(1.28), Inches(4.3), Inches(0.5),
         size=14, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)

# Main title
add_text(slide1, "BharatBuild.ai", Inches(0.5), Inches(2.2), prs.slide_width - Inches(1), Inches(1.2),
         size=68, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)

# Tagline
add_text(slide1, "Build Apps by Chatting with AI", Inches(0.5), Inches(3.3), prs.slide_width - Inches(1), Inches(0.7),
         size=28, color=GRADIENT_START, bold=False, align=PP_ALIGN.CENTER)

# Subtitle
add_text(slide1, "India's #1 AI Platform for Students, Developers & Colleges",
         Inches(0.5), Inches(4.1), prs.slide_width - Inches(1), Inches(0.5),
         size=16, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# Stats cards row
stats = [
    ("95,000+", "Students", GRADIENT_START),
    ("190+", "Companies", RGBColor(99, 102, 241)),
    ("85%", "Cost Savings", ACCENT_GREEN),
    ("< 50ms", "Latency", ACCENT_PURPLE)
]

for i, (num, label, color) in enumerate(stats):
    left = Inches(0.8 + i * 3.1)

    # Stat card
    stat_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(5), Inches(2.9), Inches(1.3))
    stat_card.fill.solid()
    stat_card.fill.fore_color.rgb = BG_CARD
    stat_card.line.color.rgb = color
    stat_card.line.width = Pt(2)

    add_text(slide1, num, left, Inches(5.15), Inches(2.9), Inches(0.55),
             size=28, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide1, label, left, Inches(5.7), Inches(2.9), Inches(0.4),
             size=12, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# Website URL with icon
add_text(slide1, "www.bharatbuild.ai", Inches(0.5), Inches(6.7), prs.slide_width - Inches(1), Inches(0.4),
         size=14, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# =====================================================
# SLIDE 2: THE PROBLEM - Pain Points
# =====================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide2)

# Section indicator
add_accent_shape(slide2, Inches(0.5), Inches(0.4), Inches(0.15), Inches(0.5), ACCENT_RED)
add_text(slide2, "THE CHALLENGE", Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
         size=14, color=ACCENT_RED, bold=True)

# Main title
add_text(slide2, "Why Current Tools Fail Indian Students", Inches(0.5), Inches(1), prs.slide_width - Inches(1), Inches(0.8),
         size=40, color=TEXT_WHITE, bold=True)

# Problem cards
problems = [
    ("Expensive Tools", "₹1,700/month", "International tools like Bolt.new & Cursor\ncost more than a student's monthly budget", ACCENT_RED),
    ("Code Only Output", "No Documentation", "Generate code but still spend weeks\nwriting SRS, UML & project reports", ACCENT_ORANGE),
    ("Payment Barriers", "Credit Card Required", "Most Indian students don't have\ninternational credit cards", ACCENT_YELLOW),
    ("High Latency", "200-400ms", "US-based servers mean slow response\ntimes on Indian internet", ACCENT_PURPLE)
]

for i, (title, stat, desc, color) in enumerate(problems):
    col = i % 2
    row = i // 2
    left = Inches(0.5 + col * 6.3)
    top = Inches(2.2 + row * 2.5)

    # Card
    add_card(slide2, left, top, Inches(6), Inches(2.2), border_color=color)

    # Icon circle
    add_icon_circle(slide2, left + Inches(0.3), top + Inches(0.3), Inches(0.7), color, "!")

    # Title
    add_text(slide2, title, left + Inches(1.2), top + Inches(0.3), Inches(3), Inches(0.5),
             size=20, color=TEXT_WHITE, bold=True)

    # Stat badge
    add_text(slide2, stat, left + Inches(4.2), top + Inches(0.35), Inches(1.5), Inches(0.4),
             size=12, color=color, bold=True, align=PP_ALIGN.RIGHT)

    # Description
    add_text(slide2, desc, left + Inches(0.3), top + Inches(1.1), Inches(5.4), Inches(1),
             size=14, color=TEXT_GRAY)

# =====================================================
# SLIDE 3: THE SOLUTION - What is BharatBuild
# =====================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide3)

# Section indicator
add_accent_shape(slide3, Inches(0.5), Inches(0.4), Inches(0.15), Inches(0.5), ACCENT_GREEN)
add_text(slide3, "THE SOLUTION", Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
         size=14, color=ACCENT_GREEN, bold=True)

# Main title
add_text(slide3, "BharatBuild.ai - Built for India", Inches(0.5), Inches(1), prs.slide_width - Inches(1), Inches(0.8),
         size=40, color=TEXT_WHITE, bold=True)

# Subtitle
add_text(slide3, "The only AI development platform that gives you code + complete academic documentation",
         Inches(0.5), Inches(1.8), prs.slide_width - Inches(1), Inches(0.5),
         size=18, color=GRADIENT_START)

# Main features
features = [
    ("Lightning Fast", "Generate full-stack applications\nin seconds, not hours", "⚡"),
    ("Complete Package", "Code + SRS + UML + Reports\n+ PPT + Viva Questions", "📦"),
    ("Indian Pricing", "Starting at just ₹99/month\n85% cheaper than competitors", "💰"),
    ("Local Servers", "Mumbai data center\n< 50ms latency", "🚀"),
    ("UPI Payments", "Pay with Google Pay, PhonePe,\nPaytm - No credit card needed", "📱"),
    ("Hindi Support", "Chat in Hindi, Tamil, Telugu\nand other Indian languages", "🗣️")
]

for i, (title, desc, icon) in enumerate(features):
    col = i % 3
    row = i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(2.8 + row * 2.2)

    add_card(slide3, left, top, Inches(4), Inches(1.9))

    # Icon
    add_text(slide3, icon, left + Inches(0.2), top + Inches(0.15), Inches(0.6), Inches(0.6),
             size=28, align=PP_ALIGN.CENTER)

    # Title
    add_text(slide3, title, left + Inches(0.9), top + Inches(0.2), Inches(2.9), Inches(0.5),
             size=18, color=TEXT_WHITE, bold=True)

    # Description
    add_text(slide3, desc, left + Inches(0.2), top + Inches(0.8), Inches(3.6), Inches(1),
             size=13, color=TEXT_GRAY)

# =====================================================
# SLIDE 4: HOW IT WORKS - Process Flow
# =====================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide4)

add_accent_shape(slide4, Inches(0.5), Inches(0.4), Inches(0.15), Inches(0.5), GRADIENT_START)
add_text(slide4, "HOW IT WORKS", Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
         size=14, color=GRADIENT_START, bold=True)

add_text(slide4, "From Idea to Deployed App in Minutes", Inches(0.5), Inches(1), prs.slide_width - Inches(1), Inches(0.8),
         size=40, color=TEXT_WHITE, bold=True)

# Process steps
steps = [
    ("1", "Describe", "Tell AI what you want to build\nin plain English or Hindi", GRADIENT_START),
    ("2", "Generate", "AI creates complete code,\ndocs, and diagrams instantly", RGBColor(99, 102, 241)),
    ("3", "Customize", "Edit code in VS Code-like editor\nwith live preview", ACCENT_PURPLE),
    ("4", "Deploy", "One-click deployment\nGet your live URL instantly", ACCENT_GREEN)
]

for i, (num, title, desc, color) in enumerate(steps):
    left = Inches(0.6 + i * 3.15)

    # Number circle
    circle = slide4.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1), Inches(2.2), Inches(1), Inches(1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    add_text(slide4, num, left + Inches(1), Inches(2.35), Inches(1), Inches(0.8),
             size=36, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Arrow (except last)
    if i < 3:
        arrow = slide4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(2.3), Inches(2.5), Inches(0.6), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(51, 65, 85)
        arrow.line.fill.background()

    # Title
    add_text(slide4, title, left, Inches(3.4), Inches(3), Inches(0.5),
             size=22, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Description
    add_text(slide4, desc, left, Inches(4), Inches(3), Inches(1),
             size=14, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# Example prompt box
example_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.3), Inches(10.3), Inches(1.5))
example_box.fill.solid()
example_box.fill.fore_color.rgb = BG_CARD
example_box.line.color.rgb = GRADIENT_START
example_box.line.width = Pt(2)

add_text(slide4, "Example Prompt:", Inches(1.8), Inches(5.45), Inches(2), Inches(0.4),
         size=12, color=GRADIENT_START, bold=True)
add_text(slide4, '"Build a hospital management system with patient registration, appointment booking, and doctor dashboard"',
         Inches(1.8), Inches(5.9), Inches(9.8), Inches(0.8),
         size=16, color=TEXT_WHITE)

# =====================================================
# SLIDE 5: UNIQUE VALUE - Academic Features
# =====================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide5)

add_accent_shape(slide5, Inches(0.5), Inches(0.4), Inches(0.15), Inches(0.5), ACCENT_ORANGE)
add_text(slide5, "ACADEMIC EXCELLENCE", Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
         size=14, color=ACCENT_ORANGE, bold=True)

add_text(slide5, "Complete Project, Not Just Code", Inches(0.5), Inches(1), prs.slide_width - Inches(1), Inches(0.8),
         size=40, color=TEXT_WHITE, bold=True)

add_text(slide5, "Everything you need for your B.Tech / M.Tech final year project",
         Inches(0.5), Inches(1.8), prs.slide_width - Inches(1), Inches(0.5),
         size=18, color=TEXT_GRAY)

# Comparison - Two columns
# Left: Others
add_card(slide5, Inches(0.5), Inches(2.5), Inches(5.8), Inches(4.3), border_color=ACCENT_RED)
add_text(slide5, "❌  Other Platforms", Inches(0.7), Inches(2.7), Inches(5.4), Inches(0.5),
         size=20, color=ACCENT_RED, bold=True)

other_items = [
    "Code files only",
    "No SRS document",
    "No UML diagrams",
    "No project report",
    "No presentation",
    "No viva preparation",
    "You spend 2-3 WEEKS extra"
]
for i, item in enumerate(other_items):
    color = ACCENT_RED if i == 6 else TEXT_GRAY
    bold = True if i == 6 else False
    add_text(slide5, f"✗  {item}", Inches(0.9), Inches(3.3 + i * 0.45), Inches(5), Inches(0.4),
             size=15, color=color, bold=bold)

# Right: BharatBuild
add_card(slide5, Inches(7), Inches(2.5), Inches(5.8), Inches(4.3), border_color=ACCENT_GREEN)
add_text(slide5, "✓  BharatBuild.ai", Inches(7.2), Inches(2.7), Inches(5.4), Inches(0.5),
         size=20, color=ACCENT_GREEN, bold=True)

bb_items = [
    "Complete source code",
    "SRS Document (IEEE format)",
    "UML Diagrams (Use Case, Class, ER, Sequence)",
    "Project Report (50+ pages)",
    "PowerPoint Presentation",
    "Viva Q&A (50+ questions with answers)",
    "Ready in 2-3 HOURS!"
]
for i, item in enumerate(bb_items):
    color = ACCENT_GREEN if i == 6 else TEXT_LIGHT
    bold = True if i == 6 else False
    add_text(slide5, f"✓  {item}", Inches(7.4), Inches(3.3 + i * 0.45), Inches(5.2), Inches(0.4),
             size=15, color=color, bold=bold)

# =====================================================
# SLIDE 6: COMPETITOR COMPARISON TABLE
# =====================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide6)

add_accent_shape(slide6, Inches(0.5), Inches(0.4), Inches(0.15), Inches(0.5), ACCENT_PURPLE)
add_text(slide6, "COMPARISON", Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
         size=14, color=ACCENT_PURPLE, bold=True)

add_text(slide6, "BharatBuild vs Competition", Inches(0.5), Inches(0.9), prs.slide_width - Inches(1), Inches(0.7),
         size=36, color=TEXT_WHITE, bold=True)

# Table
headers = ["Feature", "BharatBuild", "Bolt.new", "Cursor", "Replit"]
col_widths = [3.2, 2.4, 2.2, 2.2, 2.2]
start_x = 0.6

# Header row
for i, (header, width) in enumerate(zip(headers, col_widths)):
    left = start_x + sum(col_widths[:i])
    header_box = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(1.7), Inches(width), Inches(0.55))
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = GRADIENT_START if i == 1 else BG_CARD_LIGHT
    header_box.line.fill.background()
    add_text(slide6, header, Inches(left), Inches(1.8), Inches(width), Inches(0.4),
             size=14, color=TEXT_WHITE if i <= 1 else TEXT_GRAY, bold=True, align=PP_ALIGN.CENTER)

# Data rows
table_data = [
    ("Monthly Price", "₹99 - ₹499", "₹1,700", "₹1,700", "₹600"),
    ("Academic Docs", "✓ Complete", "✗ None", "✗ None", "✗ None"),
    ("SRS & UML", "✓ Auto-generated", "✗ Manual", "✗ Manual", "✗ Manual"),
    ("Viva Prep", "✓ 50+ Q&A", "✗ None", "✗ None", "✗ None"),
    ("Indian Payments", "✓ UPI/Paytm", "✗ Card only", "✗ Card only", "✗ Card only"),
    ("Hindi Support", "✓ Yes", "✗ No", "✗ No", "✗ No"),
    ("Server Location", "🇮🇳 India", "🇺🇸 USA", "🇺🇸 USA", "🇺🇸 USA"),
    ("Latency", "< 50ms", "200-400ms", "200-400ms", "200-400ms"),
]

for row_idx, row_data in enumerate(table_data):
    top = Inches(2.3 + row_idx * 0.58)
    bg_color = BG_CARD if row_idx % 2 == 0 else BG_CARD_LIGHT

    for col_idx, (cell, width) in enumerate(zip(row_data, col_widths)):
        left = start_x + sum(col_widths[:col_idx])

        cell_box = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), top, Inches(width), Inches(0.55))
        cell_box.fill.solid()
        cell_box.fill.fore_color.rgb = bg_color
        cell_box.line.color.rgb = RGBColor(51, 65, 85)
        cell_box.line.width = Pt(0.5)

        # Color coding
        if col_idx == 1 and "✓" in cell:
            text_color = ACCENT_GREEN
        elif "✗" in cell:
            text_color = ACCENT_RED
        elif col_idx == 0:
            text_color = TEXT_WHITE
        else:
            text_color = TEXT_GRAY

        add_text(slide6, cell, Inches(left), top + Inches(0.12), Inches(width), Inches(0.4),
                 size=13, color=text_color, align=PP_ALIGN.CENTER)

# Bottom highlight
highlight_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(6.8), Inches(6.3), Inches(0.5))
highlight_box.fill.solid()
highlight_box.fill.fore_color.rgb = ACCENT_GREEN
highlight_box.line.fill.background()
add_text(slide6, "Save 70-85% compared to international competitors!", Inches(3.5), Inches(6.85), Inches(6.3), Inches(0.4),
         size=16, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)

# =====================================================
# SLIDE 7: PRICING
# =====================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide7)

add_accent_shape(slide7, Inches(0.5), Inches(0.4), Inches(0.15), Inches(0.5), ACCENT_GREEN)
add_text(slide7, "PRICING", Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
         size=14, color=ACCENT_GREEN, bold=True)

add_text(slide7, "Affordable Plans for Everyone", Inches(0.5), Inches(0.9), prs.slide_width - Inches(1), Inches(0.7),
         size=36, color=TEXT_WHITE, bold=True)

add_text(slide7, "No hidden fees • Cancel anytime • Indian payment methods accepted",
         Inches(0.5), Inches(1.6), prs.slide_width - Inches(1), Inches(0.4),
         size=14, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# Pricing cards
plans = [
    ("Free", "₹0", "/month", ["5 projects/month", "Basic AI chat", "Code generation", "Community support"], False),
    ("Student", "₹99", "/month", ["Unlimited projects", "Full AI features", "SRS & UML generation", "Email support", "Viva Q&A"], False),
    ("Pro", "₹299", "/month", ["Everything in Student", "Priority AI (faster)", "Docker execution", "Project reports", "Priority support"], True),
    ("College", "₹499", "/student/year", ["Bulk licensing", "Faculty dashboard", "Plagiarism detection", "Analytics & reports", "Dedicated support"], False)
]

for i, (name, price, period, features, popular) in enumerate(plans):
    left = Inches(0.5 + i * 3.15)

    # Card with highlight for popular
    card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.2), Inches(3), Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD
    card.line.color.rgb = ACCENT_GREEN if popular else RGBColor(51, 65, 85)
    card.line.width = Pt(3) if popular else Pt(1)

    # Popular badge
    if popular:
        badge = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.5), Inches(1.95), Inches(2), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT_GREEN
        badge.line.fill.background()
        add_text(slide7, "MOST POPULAR", left + Inches(0.5), Inches(1.98), Inches(2), Inches(0.3),
                 size=11, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)

    # Plan name
    add_text(slide7, name, left, Inches(2.5), Inches(3), Inches(0.5),
             size=22, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Price
    add_text(slide7, price, left, Inches(3), Inches(3), Inches(0.7),
             size=40, color=GRADIENT_START, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide7, period, left, Inches(3.6), Inches(3), Inches(0.3),
             size=12, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    # Features
    for j, feature in enumerate(features):
        add_text(slide7, f"✓ {feature}", left + Inches(0.2), Inches(4.1 + j * 0.45), Inches(2.6), Inches(0.4),
                 size=12, color=TEXT_GRAY)

# =====================================================
# SLIDE 8: TARGET AUDIENCE
# =====================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide8)

add_accent_shape(slide8, Inches(0.5), Inches(0.4), Inches(0.15), Inches(0.5), GRADIENT_START)
add_text(slide8, "WHO IS IT FOR", Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
         size=14, color=GRADIENT_START, bold=True)

add_text(slide8, "Built for Every Builder", Inches(0.5), Inches(0.9), prs.slide_width - Inches(1), Inches(0.7),
         size=36, color=TEXT_WHITE, bold=True)

audiences = [
    ("🎓", "Students", "Complete final year projects\nwith all documentation", "Save 2-3 weeks per project", GRADIENT_START),
    ("👨‍💻", "Developers", "Rapid prototyping &\nproduction-ready apps", "5x faster development", RGBColor(99, 102, 241)),
    ("🚀", "Founders", "Build MVPs without\nhiring developers", "Launch in days, not months", ACCENT_PURPLE),
    ("🏫", "Colleges", "Track student progress &\ndetect plagiarism", "Manage 1000+ students", ACCENT_GREEN)
]

for i, (icon, title, desc, benefit, color) in enumerate(audiences):
    left = Inches(0.5 + i * 3.15)

    add_card(slide8, left, Inches(1.9), Inches(3), Inches(4.5), border_color=color)

    # Icon
    add_text(slide8, icon, left, Inches(2.2), Inches(3), Inches(0.8),
             size=48, align=PP_ALIGN.CENTER)

    # Title
    add_text(slide8, title, left, Inches(3.1), Inches(3), Inches(0.5),
             size=22, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Description
    add_text(slide8, desc, left + Inches(0.2), Inches(3.7), Inches(2.6), Inches(1),
             size=14, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    # Benefit badge
    benefit_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.3), Inches(5.6), Inches(2.4), Inches(0.5))
    benefit_box.fill.solid()
    benefit_box.fill.fore_color.rgb = color
    benefit_box.line.fill.background()
    add_text(slide8, benefit, left + Inches(0.3), Inches(5.65), Inches(2.4), Inches(0.4),
             size=11, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)

# =====================================================
# SLIDE 9: TECHNOLOGY STACK
# =====================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide9)

add_accent_shape(slide9, Inches(0.5), Inches(0.4), Inches(0.15), Inches(0.5), ACCENT_PURPLE)
add_text(slide9, "TECHNOLOGY", Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
         size=14, color=ACCENT_PURPLE, bold=True)

add_text(slide9, "Enterprise-Grade Technology", Inches(0.5), Inches(0.9), prs.slide_width - Inches(1), Inches(0.7),
         size=36, color=TEXT_WHITE, bold=True)

# Tech categories
tech_sections = [
    ("Frontend", ["Next.js 14 + React 18", "Monaco Editor (VS Code)", "TypeScript", "Tailwind CSS"], GRADIENT_START),
    ("Backend", ["FastAPI (Python 3.11)", "Claude 3.5 Sonnet AI", "PostgreSQL + Redis", "Docker Containers"], RGBColor(99, 102, 241)),
    ("Infrastructure", ["Indian Servers (Mumbai)", "< 50ms Latency", "Auto-scaling", "99.9% Uptime"], ACCENT_PURPLE),
    ("Security", ["JWT Authentication", "Isolated Sandboxes", "Data Encryption", "GDPR Compliant"], ACCENT_GREEN)
]

for i, (title, items, color) in enumerate(tech_sections):
    col = i % 2
    row = i // 2
    left = Inches(0.5 + col * 6.4)
    top = Inches(1.8 + row * 2.7)

    add_card(slide9, left, top, Inches(6.1), Inches(2.4))

    # Color bar
    bar = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.15), Inches(2.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    # Title
    add_text(slide9, title, left + Inches(0.4), top + Inches(0.2), Inches(3), Inches(0.5),
             size=20, color=color, bold=True)

    # Items
    for j, item in enumerate(items):
        add_text(slide9, f"• {item}", left + Inches(0.4), top + Inches(0.75 + j * 0.4), Inches(5.5), Inches(0.4),
                 size=14, color=TEXT_GRAY)

# =====================================================
# SLIDE 10: CAMPUS DRIVE FEATURE
# =====================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide10)

add_accent_shape(slide10, Inches(0.5), Inches(0.4), Inches(0.15), Inches(0.5), ACCENT_ORANGE)
add_text(slide10, "CAMPUS RECRUITMENT", Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
         size=14, color=ACCENT_ORANGE, bold=True)

add_text(slide10, "Integrated Assessment Platform", Inches(0.5), Inches(0.9), prs.slide_width - Inches(1), Inches(0.7),
         size=36, color=TEXT_WHITE, bold=True)

add_text(slide10, "Conduct campus placement drives for 300+ concurrent students",
         Inches(0.5), Inches(1.6), prs.slide_width - Inches(1), Inches(0.4),
         size=16, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# Features
campus_features = [
    ("📝", "Multi-Category Quiz", "Logical, Technical, AI/ML,\nEnglish & Coding MCQs"),
    ("⚡", "Auto-Evaluation", "Instant results with\nsection-wise scores"),
    ("💾", "Auto-Save & Resume", "Progress saved automatically\nResume if disconnected"),
    ("📊", "Admin Dashboard", "Track registrations &\nview detailed analytics"),
    ("🔀", "Randomization", "Unique question order\nfor each student"),
    ("🔒", "Secure & Scalable", "300+ concurrent users\nwith no performance issues")
]

for i, (icon, title, desc) in enumerate(campus_features):
    col = i % 3
    row = i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(2.3 + row * 2.4)

    add_card(slide10, left, top, Inches(4), Inches(2.1))

    add_text(slide10, icon, left + Inches(0.2), top + Inches(0.2), Inches(0.6), Inches(0.6), size=28)
    add_text(slide10, title, left + Inches(0.9), top + Inches(0.25), Inches(2.9), Inches(0.5),
             size=16, color=TEXT_WHITE, bold=True)
    add_text(slide10, desc, left + Inches(0.2), top + Inches(0.9), Inches(3.6), Inches(1),
             size=13, color=TEXT_GRAY)

# =====================================================
# SLIDE 11: TESTIMONIALS
# =====================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide11)

add_accent_shape(slide11, Inches(0.5), Inches(0.4), Inches(0.15), Inches(0.5), ACCENT_YELLOW)
add_text(slide11, "TESTIMONIALS", Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
         size=14, color=ACCENT_YELLOW, bold=True)

add_text(slide11, "Loved by Students & Educators", Inches(0.5), Inches(0.9), prs.slide_width - Inches(1), Inches(0.7),
         size=36, color=TEXT_WHITE, bold=True)

testimonials = [
    ('"BharatBuild saved my final year project! Got complete code, SRS, UML diagrams, and even viva questions. My guide was impressed!"',
     "Priya Sharma", "B.Tech Student, VIT Vellore", "⭐⭐⭐⭐⭐"),
    ('"I compared it with Bolt.new - same features but 17x cheaper. Perfect for student budget. The Hindi support is a game changer!"',
     "Rahul Kumar", "CS Student, IIT Delhi", "⭐⭐⭐⭐⭐"),
    ('"Managing 60 students\' projects was impossible before. Now I can track everyone\'s progress from one dashboard. Highly recommend!"',
     "Dr. Anitha Reddy", "Professor, Anna University", "⭐⭐⭐⭐⭐")
]

for i, (quote, name, role, stars) in enumerate(testimonials):
    top = Inches(1.7 + i * 1.85)

    add_card(slide11, Inches(0.8), top, Inches(11.7), Inches(1.6))

    # Quote
    add_text(slide11, quote, Inches(1.1), top + Inches(0.15), Inches(9.5), Inches(0.9),
             size=15, color=TEXT_LIGHT)

    # Attribution
    add_text(slide11, f"{name}  •  {role}", Inches(1.1), top + Inches(1.05), Inches(7), Inches(0.4),
             size=13, color=GRADIENT_START, bold=True)

    # Stars
    add_text(slide11, stars, Inches(10), top + Inches(1.05), Inches(2.3), Inches(0.4),
             size=14, color=ACCENT_YELLOW, align=PP_ALIGN.RIGHT)

# =====================================================
# SLIDE 12: MARKET OPPORTUNITY
# =====================================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide12)

add_accent_shape(slide12, Inches(0.5), Inches(0.4), Inches(0.15), Inches(0.5), ACCENT_GREEN)
add_text(slide12, "MARKET OPPORTUNITY", Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
         size=14, color=ACCENT_GREEN, bold=True)

add_text(slide12, "Massive Growth Potential", Inches(0.5), Inches(0.9), prs.slide_width - Inches(1), Inches(0.7),
         size=36, color=TEXT_WHITE, bold=True)

# TAM
add_text(slide12, "Total Addressable Market", Inches(0.5), Inches(1.7), Inches(4), Inches(0.4),
         size=16, color=TEXT_GRAY)
add_text(slide12, "₹2,195 Crore", Inches(0.5), Inches(2.1), Inches(4), Inches(0.7),
         size=48, color=ACCENT_GREEN, bold=True)

# Market segments
segments = [
    ("3,500+", "Engineering Colleges", "₹1,750 Cr"),
    ("10M+", "Engineering Students", "Direct B2C"),
    ("130+", "Deemed Universities", "₹195 Cr"),
    ("50+", "State Universities", "₹250 Cr")
]

for i, (num, segment, value) in enumerate(segments):
    left = Inches(5.5 + (i % 2) * 4)
    top = Inches(1.7 + (i // 2) * 1.8)

    add_card(slide12, left, top, Inches(3.7), Inches(1.5))
    add_text(slide12, num, left + Inches(0.2), top + Inches(0.2), Inches(3.3), Inches(0.6),
             size=32, color=GRADIENT_START, bold=True)
    add_text(slide12, segment, left + Inches(0.2), top + Inches(0.8), Inches(3.3), Inches(0.35),
             size=14, color=TEXT_WHITE)
    add_text(slide12, value, left + Inches(0.2), top + Inches(1.1), Inches(3.3), Inches(0.3),
             size=12, color=TEXT_GRAY)

# Growth projection
add_text(slide12, "Growth Trajectory", Inches(0.5), Inches(4.4), Inches(4), Inches(0.4),
         size=18, color=TEXT_WHITE, bold=True)

projections = [
    ("Year 1", "20 colleges", "₹80L ARR"),
    ("Year 2", "80 colleges", "₹4 Cr ARR"),
    ("Year 3", "200 colleges", "₹12 Cr ARR")
]

for i, (year, colleges, revenue) in enumerate(projections):
    left = Inches(0.5 + i * 2.2)
    add_text(slide12, year, left, Inches(4.9), Inches(2), Inches(0.4),
             size=14, color=TEXT_GRAY)
    add_text(slide12, colleges, left, Inches(5.3), Inches(2), Inches(0.4),
             size=16, color=TEXT_WHITE, bold=True)
    add_text(slide12, revenue, left, Inches(5.7), Inches(2), Inches(0.4),
             size=14, color=ACCENT_GREEN)

# Break-even
add_text(slide12, "Break-even: Year 2 (~80 colleges)", Inches(0.5), Inches(6.5), Inches(6), Inches(0.4),
         size=16, color=ACCENT_ORANGE, bold=True)

# =====================================================
# SLIDE 13: CALL TO ACTION (Clean Design)
# =====================================================
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide13)

# Top accent bar
cta_bar = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
cta_bar.fill.solid()
cta_bar.fill.fore_color.rgb = ACCENT_GREEN
cta_bar.line.fill.background()

# Main CTA heading
add_text(slide13, "Ready to Build Your Next Project?", Inches(0.5), Inches(1.5), prs.slide_width - Inches(1), Inches(0.9),
         size=44, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide13, "Join 95,000+ students & developers already building with AI",
         Inches(0.5), Inches(2.5), prs.slide_width - Inches(1), Inches(0.5),
         size=20, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# Feature highlights in a row
cta_features = [
    ("Free to Start", "No credit card required"),
    ("Indian Payments", "UPI, Paytm, PhonePe"),
    ("Instant Access", "Start building in seconds")
]

for i, (title, desc) in enumerate(cta_features):
    left = Inches(1 + i * 4)

    feature_box = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.3), Inches(3.5), Inches(1.2))
    feature_box.fill.solid()
    feature_box.fill.fore_color.rgb = BG_CARD
    feature_box.line.color.rgb = GRADIENT_START
    feature_box.line.width = Pt(1)

    add_text(slide13, title, left, Inches(3.45), Inches(3.5), Inches(0.45),
             size=16, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide13, desc, left, Inches(3.95), Inches(3.5), Inches(0.4),
             size=12, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# CTA Button - larger and more prominent
cta_btn = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(4.9), Inches(5.3), Inches(0.85))
cta_btn.fill.solid()
cta_btn.fill.fore_color.rgb = ACCENT_GREEN
cta_btn.line.fill.background()
add_text(slide13, "Start Building for FREE", Inches(4), Inches(5.1), Inches(5.3), Inches(0.55),
         size=22, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)

# Website URL
add_text(slide13, "www.bharatbuild.ai", Inches(0.5), Inches(6), prs.slide_width - Inches(1), Inches(0.5),
         size=28, color=GRADIENT_START, bold=True, align=PP_ALIGN.CENTER)

# Contact info
add_text(slide13, "support@bharatbuild.ai  |  Made in India",
         Inches(0.5), Inches(6.7), prs.slide_width - Inches(1), Inches(0.4),
         size=14, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# =====================================================
# SLIDE 14: THANK YOU
# =====================================================
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide14)

add_text(slide14, "Thank You", Inches(0.5), Inches(2.5), prs.slide_width - Inches(1), Inches(1),
         size=72, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide14, "Questions?", Inches(0.5), Inches(3.8), prs.slide_width - Inches(1), Inches(0.6),
         size=28, color=GRADIENT_START, align=PP_ALIGN.CENTER)

# Contact details
add_text(slide14, "🌐  www.bharatbuild.ai", Inches(0.5), Inches(5), prs.slide_width - Inches(1), Inches(0.5),
         size=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide14, "📧  support@bharatbuild.ai", Inches(0.5), Inches(5.5), prs.slide_width - Inches(1), Inches(0.5),
         size=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide14, "Build Apps by Chatting with AI", Inches(0.5), Inches(6.5), prs.slide_width - Inches(1), Inches(0.5),
         size=16, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

# =====================================================
# SAVE PRESENTATION
# =====================================================
output_path = os.path.join(os.path.dirname(__file__), "BharatBuild_Presentation_Final.pptx")
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print("\nSlide Contents:")
print("  1. Title - Hero Section")
print("  2. The Problem - Pain Points")
print("  3. The Solution - BharatBuild Features")
print("  4. How It Works - Process Flow")
print("  5. Academic Excellence - Complete Package")
print("  6. Competitor Comparison Table")
print("  7. Pricing Plans")
print("  8. Target Audience")
print("  9. Technology Stack")
print(" 10. Campus Recruitment Feature")
print(" 11. Testimonials")
print(" 12. Market Opportunity")
print(" 13. Call to Action")
print(" 14. Thank You")
