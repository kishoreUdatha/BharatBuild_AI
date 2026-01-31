"""
BharatBuild.ai - Comprehensive Investor Pitch Deck
Professional presentation for fundraising with detailed information
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Capture website screenshots
def capture_screenshots():
    """Capture screenshots from bharatbuild.ai"""
    screenshots = {}
    script_dir = os.path.dirname(os.path.abspath(__file__))

    try:
        from playwright.sync_api import sync_playwright

        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={'width': 1920, 'height': 1080})

            pages_to_capture = [
                ('homepage', 'https://bharatbuild.ai/', 'Homepage'),
                ('bolt', 'https://bharatbuild.ai/bolt', 'AI Editor'),
                ('pricing', 'https://bharatbuild.ai/pricing', 'Pricing'),
                ('research', 'https://bharatbuild.ai/research-paper', 'Research Paper'),
            ]

            for i, (key, url, name) in enumerate(pages_to_capture, 1):
                print(f"{i}. Capturing {name}...")
                page.goto(url, wait_until="networkidle", timeout=60000)
                page.wait_for_timeout(3000)
                screenshots[key] = os.path.join(script_dir, f"screenshot_{key}.png")
                page.screenshot(path=screenshots[key])
                print(f"   Saved: {screenshots[key]}")

            # Capture features section (scrolled)
            print("5. Capturing Features section...")
            page.goto("https://bharatbuild.ai/", wait_until="networkidle", timeout=60000)
            page.evaluate("window.scrollBy(0, 800)")
            page.wait_for_timeout(2000)
            screenshots['features'] = os.path.join(script_dir, "screenshot_features.png")
            page.screenshot(path=screenshots['features'])
            print(f"   Saved: {screenshots['features']}")

            browser.close()
            print(f"All {len(screenshots)} screenshots captured!")

    except Exception as e:
        print(f"Screenshot capture error: {e}")

    return screenshots

print("Capturing website screenshots...")
website_screenshots = capture_screenshots()

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
BG_DARK = RGBColor(10, 10, 20)
BG_CARD = RGBColor(20, 25, 40)
BG_CARD_LIGHT = RGBColor(35, 40, 60)
TEXT_WHITE = RGBColor(255, 255, 255)
TEXT_GRAY = RGBColor(160, 170, 190)
TEXT_LIGHT = RGBColor(220, 225, 235)
PRIMARY = RGBColor(0, 200, 220)
ACCENT_GREEN = RGBColor(40, 200, 100)
ACCENT_RED = RGBColor(240, 80, 80)
ACCENT_ORANGE = RGBColor(255, 160, 50)
ACCENT_PURPLE = RGBColor(140, 100, 240)
ACCENT_BLUE = RGBColor(80, 130, 250)

def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_text(slide, text, left, top, width, height, size=18, color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_card(slide, left, top, width, height, border_color=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(2)
    else:
        card.line.color.rgb = RGBColor(50, 60, 80)
        card.line.width = Pt(1)
    return card

def add_slide_number(slide, num, total):
    add_text(slide, f"{num}/{total}", Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3),
             size=10, color=TEXT_GRAY, align=PP_ALIGN.RIGHT)

TOTAL_SLIDES = 16

# =====================================================
# SLIDE 1: TITLE / COVER
# =====================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1)

bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = PRIMARY
bar.line.fill.background()

add_text(slide1, "BharatBuild.ai", Inches(0.5), Inches(2.2), prs.slide_width - Inches(1), Inches(1),
         size=72, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide1, "India's First AI-Powered Academic Project Platform", Inches(0.5), Inches(3.4), prs.slide_width - Inches(1), Inches(0.6),
         size=28, color=PRIMARY, align=PP_ALIGN.CENTER)

add_text(slide1, "Empowering 10M+ engineering students to build production-ready applications\nwith complete academic documentation using AI", Inches(1), Inches(4.3), prs.slide_width - Inches(2), Inches(0.9),
         size=18, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# Key stats
stats = [("Rs 4,999", "Per Project"), ("10M+", "Target Students"), ("Rs 13.5 Cr", "Year 3 Revenue")]
for i, (num, label) in enumerate(stats):
    left = Inches(3 + i * 2.8)
    add_card(slide1, left, Inches(5.3), Inches(2.5), Inches(1.1))
    add_text(slide1, num, left, Inches(5.45), Inches(2.5), Inches(0.5), size=20, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide1, label, left, Inches(5.95), Inches(2.5), Inches(0.3), size=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(6.6), Inches(2.3), Inches(0.45))
badge.fill.solid()
badge.fill.fore_color.rgb = ACCENT_GREEN
badge.line.fill.background()
add_text(slide1, "Seed Round", Inches(5.5), Inches(6.65), Inches(2.3), Inches(0.35), size=14, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)

add_text(slide1, "www.bharatbuild.ai | Confidential", Inches(0.5), Inches(7.1), prs.slide_width - Inches(1), Inches(0.3),
         size=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
add_slide_number(slide1, 1, TOTAL_SLIDES)

# =====================================================
# SLIDE 2: THE PROBLEM
# =====================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)

add_text(slide2, "The Problem", Inches(0.5), Inches(0.25), Inches(5), Inches(0.5), size=32, color=TEXT_WHITE, bold=True)
add_text(slide2, "Indian engineering students face critical challenges", Inches(0.5), Inches(0.75), Inches(12), Inches(0.4), size=14, color=TEXT_GRAY)

problems = [
    ("10M+", "Students Struggle", "Final year projects need code + SRS + UML + reports + PPT + viva. AI tools only generate code.", ACCENT_RED),
    ("Rs 1,700/mo", "Unaffordable Tools", "Bolt.new, Cursor cost $20+/month. More than most students' monthly pocket money.", ACCENT_ORANGE),
    ("70%+", "Payment Barrier", "International tools need credit cards. Most Indian students don't have access.", ACCENT_PURPLE),
    ("2-3 Weeks", "Wasted Time", "Students spend weeks on documentation instead of learning. Many copy and fail viva.", ACCENT_BLUE)
]

for i, (stat, title, desc, color) in enumerate(problems):
    col = i % 2
    row = i // 2
    left = Inches(0.4 + col * 6.5)
    top = Inches(1.25 + row * 3.0)

    add_card(slide2, left, top, Inches(6.3), Inches(2.8), border_color=color)
    add_text(slide2, stat, left + Inches(0.3), top + Inches(0.25), Inches(3), Inches(0.6), size=36, color=color, bold=True)
    add_text(slide2, title, left + Inches(0.3), top + Inches(0.9), Inches(5.7), Inches(0.4), size=16, color=TEXT_WHITE, bold=True)
    add_text(slide2, desc, left + Inches(0.3), top + Inches(1.4), Inches(5.7), Inches(1.2), size=12, color=TEXT_GRAY)

add_slide_number(slide2, 2, TOTAL_SLIDES)

# =====================================================
# SLIDE 3: OUR SOLUTION
# =====================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3)

add_text(slide3, "Our Solution", Inches(0.5), Inches(0.3), Inches(5), Inches(0.6), size=36, color=TEXT_WHITE, bold=True)
add_text(slide3, "BharatBuild.ai - Complete AI Development Platform for Indian Students", Inches(0.5), Inches(0.95), Inches(12), Inches(0.5), size=18, color=PRIMARY)

solutions = [
    ("AI Code Generation", "Build full-stack applications by chatting with AI. Support for React, Node.js, Python, and more. Live preview in browser.", PRIMARY),
    ("Complete Documentation", "Auto-generate SRS, UML diagrams (Class, Sequence, ER, Use Case), project reports, PPT presentations.", ACCENT_GREEN),
    ("Viva Preparation", "50+ viva Q&A generated based on your project. Understand your own code better.", ACCENT_PURPLE),
    ("One-Time Payment", "Rs 4,999 per project. No monthly subscription. Pay once, get everything.", ACCENT_ORANGE),
    ("Indian Payments", "UPI, Paytm, PhonePe, Net Banking. No credit card required. Instant activation.", ACCENT_BLUE),
    ("Low Latency", "Servers in India (Mumbai). <50ms response time. 4x faster than US-based competitors.", ACCENT_GREEN)
]

for i, (title, desc, color) in enumerate(solutions):
    col = i % 3
    row = i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(1.6 + row * 2.7)

    add_card(slide3, left, top, Inches(4), Inches(2.5))
    add_text(slide3, title, left + Inches(0.2), top + Inches(0.2), Inches(3.6), Inches(0.5), size=16, color=color, bold=True)
    add_text(slide3, desc, left + Inches(0.2), top + Inches(0.75), Inches(3.6), Inches(1.6), size=13, color=TEXT_GRAY)

add_slide_number(slide3, 3, TOTAL_SLIDES)

# =====================================================
# SLIDE 4: PRODUCT SCREENSHOTS
# =====================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4)

add_text(slide4, "Product Overview", Inches(0.5), Inches(0.3), Inches(8), Inches(0.5), size=36, color=TEXT_WHITE, bold=True)

screenshot_config = [
    ('homepage', 'Homepage', PRIMARY, 0, 0),
    ('bolt', 'AI Code Editor', ACCENT_GREEN, 1, 0),
    ('pricing', 'Pricing (Rs 4,999)', ACCENT_ORANGE, 2, 0),
    ('features', 'Platform Features', ACCENT_BLUE, 0, 1),
    ('research', 'Auto Documentation', ACCENT_PURPLE, 1, 1),
]

img_width, img_height = 4.0, 2.4
start_left, start_top = 0.5, 0.85
col_gap, row_gap = 4.2, 2.75

for key, label, color, col, row in screenshot_config:
    left = Inches(start_left + col * col_gap)
    top = Inches(start_top + row * row_gap)

    if key in website_screenshots and os.path.exists(website_screenshots[key]):
        slide4.shapes.add_picture(website_screenshots[key], left, top, width=Inches(img_width), height=Inches(img_height))
    else:
        placeholder = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(img_width), Inches(img_height))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = BG_CARD
        placeholder.line.color.rgb = color
        placeholder.line.width = Pt(2)

    add_text(slide4, label, left, top + Inches(img_height + 0.05), Inches(img_width), Inches(0.3),
             size=11, color=color, bold=True, align=PP_ALIGN.CENTER)

add_slide_number(slide4, 4, TOTAL_SLIDES)

# =====================================================
# SLIDE 5: HOW IT WORKS
# =====================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5)

add_text(slide5, "How It Works", Inches(0.5), Inches(0.3), Inches(5), Inches(0.6), size=36, color=TEXT_WHITE, bold=True)
add_text(slide5, "From idea to complete project in 4 simple steps", Inches(0.5), Inches(0.95), Inches(12), Inches(0.5), size=16, color=TEXT_GRAY)

steps = [
    ("1", "Describe Your Project", "Tell AI what you want to build.\n'Build a hospital management\nsystem with patient records,\nappointments, and billing'", PRIMARY),
    ("2", "AI Generates Code", "Full-stack code generated in\nminutes. React frontend,\nNode.js backend, database\nschema - all production ready.", ACCENT_GREEN),
    ("3", "Preview & Edit", "See live preview in browser.\nEdit code with Monaco editor\n(VS Code experience).\nIterate with AI assistance.", ACCENT_ORANGE),
    ("4", "Download Everything", "Get complete package:\n- Source code\n- SRS document\n- UML diagrams\n- Project report\n- PPT presentation\n- 50+ Viva Q&A", ACCENT_PURPLE)
]

for i, (num, title, desc, color) in enumerate(steps):
    left = Inches(0.5 + i * 3.2)

    # Step circle
    circle = slide5.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.1), Inches(1.6), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(slide5, num, left + Inches(1.1), Inches(1.7), Inches(0.7), Inches(0.5), size=24, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)

    # Arrow (except last)
    if i < 3:
        arrow = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(2.5), Inches(1.8), Inches(0.6), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(60, 70, 90)
        arrow.line.fill.background()

    add_card(slide5, left, Inches(2.5), Inches(3), Inches(4.2), border_color=color)
    add_text(slide5, title, left + Inches(0.15), Inches(2.7), Inches(2.7), Inches(0.5), size=15, color=color, bold=True)
    add_text(slide5, desc, left + Inches(0.15), Inches(3.2), Inches(2.7), Inches(3.3), size=12, color=TEXT_GRAY)

add_slide_number(slide5, 5, TOTAL_SLIDES)

# =====================================================
# SLIDE 6: MARKET OPPORTUNITY
# =====================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6)

add_text(slide6, "Market Opportunity", Inches(0.5), Inches(0.3), Inches(5), Inches(0.6), size=36, color=TEXT_WHITE, bold=True)

markets = [
    ("TAM", "Rs 2,195 Cr", "Total Addressable Market", "All engineering colleges (3,500+)\n+ All students (10M+) in India", ACCENT_BLUE),
    ("SAM", "Rs 500 Cr", "Serviceable Market", "Private engineering colleges\n+ Direct B2C students\nwho can afford Rs 4,999", ACCENT_PURPLE),
    ("SOM", "Rs 13.5 Cr", "Year 3 Target", "100 colleges (B2B)\n+ 20,000 paid users (B2C)", ACCENT_GREEN)
]

for i, (label, value, title, desc, color) in enumerate(markets):
    left = Inches(0.5 + i * 4.2)
    add_card(slide6, left, Inches(1.1), Inches(4), Inches(2.6), border_color=color)
    add_text(slide6, label, left + Inches(0.25), Inches(1.25), Inches(1.5), Inches(0.4), size=14, color=color, bold=True)
    add_text(slide6, value, left + Inches(0.25), Inches(1.65), Inches(3.5), Inches(0.6), size=32, color=TEXT_WHITE, bold=True)
    add_text(slide6, title, left + Inches(0.25), Inches(2.3), Inches(3.5), Inches(0.35), size=12, color=color)
    add_text(slide6, desc, left + Inches(0.25), Inches(2.7), Inches(3.5), Inches(0.9), size=11, color=TEXT_GRAY)

add_text(slide6, "Why This Market?", Inches(0.5), Inches(3.9), Inches(4), Inches(0.4), size=18, color=TEXT_WHITE, bold=True)

reasons = [
    ("3,500+", "Engineering Colleges", "in India"),
    ("10M+", "Engineering Students", "graduate annually"),
    ("Rs 4,999", "Affordable Price Point", "for Indian students"),
    ("NEP 2020", "Policy Push", "for practical learning")
]

for i, (num, label, sub) in enumerate(reasons):
    left = Inches(0.5 + i * 3.2)
    add_text(slide6, num, left, Inches(4.4), Inches(3), Inches(0.5), size=28, color=PRIMARY, bold=True)
    add_text(slide6, label, left, Inches(4.95), Inches(3), Inches(0.35), size=13, color=TEXT_WHITE)
    add_text(slide6, sub, left, Inches(5.3), Inches(3), Inches(0.3), size=11, color=TEXT_GRAY)

add_text(slide6, "Growth Drivers: NEP 2020 practical learning focus | Rising EdTech adoption (40% YoY) | AI tools market explosion | Increasing internet penetration in Tier 2/3 cities",
         Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.8), size=12, color=TEXT_GRAY)

add_slide_number(slide6, 6, TOTAL_SLIDES)

# =====================================================
# SLIDE 7: BUSINESS MODEL
# =====================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7)

add_text(slide7, "Business Model", Inches(0.5), Inches(0.3), Inches(5), Inches(0.6), size=36, color=TEXT_WHITE, bold=True)
add_text(slide7, "Dual Revenue Streams: B2C Per-Project + B2B College Licensing", Inches(0.5), Inches(0.9), Inches(12), Inches(0.4), size=16, color=PRIMARY)

# B2C Model
add_card(slide7, Inches(0.5), Inches(1.4), Inches(6.2), Inches(2.8))
add_text(slide7, "B2C - Per Project (One-Time Payment)", Inches(0.7), Inches(1.55), Inches(5.8), Inches(0.4), size=15, color=ACCENT_GREEN, bold=True)

b2c_plans = [
    ("Free Tier", "Rs 0", "3 preview files only\n7-day validity\nNo downloads"),
    ("Premium", "Rs 4,999", "Complete project\nAll source code\nSRS + UML + Report\nPPT + Viva Q&A\nLifetime access")
]

for i, (name, price, features) in enumerate(b2c_plans):
    left = Inches(0.9 + i * 3)
    add_text(slide7, name, left, Inches(2.05), Inches(2.7), Inches(0.35), size=14, color=TEXT_WHITE, bold=True)
    add_text(slide7, price, left, Inches(2.4), Inches(2.7), Inches(0.45), size=22, color=PRIMARY, bold=True)
    add_text(slide7, features, left, Inches(2.9), Inches(2.7), Inches(1.2), size=11, color=TEXT_GRAY)

# B2B Model
add_card(slide7, Inches(6.9), Inches(1.4), Inches(6.2), Inches(2.8))
add_text(slide7, "B2B - College Licensing (Annual)", Inches(7.1), Inches(1.55), Inches(5.8), Inches(0.4), size=15, color=ACCENT_BLUE, bold=True)

b2b_plans = [
    ("Starter", "Rs 2L/yr", "200 students"),
    ("Standard", "Rs 3L/yr", "500 students"),
    ("Premium", "Rs 3.5L/yr", "1000 students")
]

for i, (name, price, students) in enumerate(b2b_plans):
    left = Inches(7.3 + i * 1.9)
    add_text(slide7, name, left, Inches(2.05), Inches(1.7), Inches(0.35), size=13, color=TEXT_WHITE, bold=True)
    add_text(slide7, price, left, Inches(2.4), Inches(1.7), Inches(0.4), size=16, color=PRIMARY, bold=True)
    add_text(slide7, students, left, Inches(2.85), Inches(1.7), Inches(0.3), size=11, color=TEXT_GRAY)

add_text(slide7, "B2B Includes: Admin dashboard, usage analytics, bulk student management, priority support, custom branding",
         Inches(7.1), Inches(3.3), Inches(5.8), Inches(0.7), size=10, color=TEXT_GRAY)

# Unit Economics
add_text(slide7, "Unit Economics", Inches(0.5), Inches(4.4), Inches(4), Inches(0.4), size=18, color=TEXT_WHITE, bold=True)

metrics = [
    ("B2C Revenue", "Rs 4,999", "per project"),
    ("B2B Revenue", "Rs 2-3.5L", "per college/year"),
    ("Gross Margin", "70%", "after AI & cloud"),
    ("CAC (B2C)", "Rs 500", "via social/referral"),
    ("LTV:CAC", "10:1", "strong unit economics")
]

for i, (label, value, note) in enumerate(metrics):
    left = Inches(0.5 + i * 2.5)
    add_card(slide7, left, Inches(4.9), Inches(2.3), Inches(1.1))
    add_text(slide7, label, left + Inches(0.1), Inches(5.0), Inches(2.1), Inches(0.25), size=10, color=TEXT_GRAY)
    add_text(slide7, value, left + Inches(0.1), Inches(5.25), Inches(2.1), Inches(0.35), size=16, color=ACCENT_GREEN, bold=True)
    add_text(slide7, note, left + Inches(0.1), Inches(5.6), Inches(2.1), Inches(0.25), size=9, color=TEXT_GRAY)

add_slide_number(slide7, 7, TOTAL_SLIDES)

# =====================================================
# SLIDE 8: GO-TO-MARKET STRATEGY
# =====================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8)

add_text(slide8, "Go-to-Market Strategy", Inches(0.5), Inches(0.3), Inches(6), Inches(0.6), size=36, color=TEXT_WHITE, bold=True)

# B2C GTM
add_text(slide8, "B2C Strategy - Direct to Students", Inches(0.5), Inches(1.0), Inches(6), Inches(0.4), size=18, color=ACCENT_GREEN, bold=True)

b2c_channels = [
    ("YouTube & Reels", "Tutorial videos showcasing\nproject builds. Target: 1M views/month"),
    ("College Ambassadors", "Student ambassadors in\n100+ colleges. 10% referral commission"),
    ("SEO & Content", "Rank for 'final year project'\n'SRS generator' keywords"),
    ("WhatsApp Groups", "Direct outreach to student\ncommunities. Viral sharing")
]

for i, (channel, desc) in enumerate(b2c_channels):
    left = Inches(0.5 + i * 3.2)
    add_card(slide8, left, Inches(1.5), Inches(3), Inches(1.8))
    add_text(slide8, channel, left + Inches(0.15), Inches(1.65), Inches(2.7), Inches(0.35), size=12, color=PRIMARY, bold=True)
    add_text(slide8, desc, left + Inches(0.15), Inches(2.05), Inches(2.7), Inches(1.1), size=10, color=TEXT_GRAY)

# B2B GTM
add_text(slide8, "B2B Strategy - College Partnerships", Inches(0.5), Inches(3.5), Inches(6), Inches(0.4), size=18, color=ACCENT_BLUE, bold=True)

b2b_channels = [
    ("Direct Sales", "Dedicated sales team targeting\nHODs and Principals. 6-month cycle"),
    ("EdTech Conferences", "Presence at FICCI, CII education\nevents. Demo booths"),
    ("Pilot Programs", "Free 3-month pilots for\ntop colleges. Convert to paid"),
    ("Government Tenders", "AICTE, state govt partnerships\nfor skill development")
]

for i, (channel, desc) in enumerate(b2b_channels):
    left = Inches(0.5 + i * 3.2)
    add_card(slide8, left, Inches(4.0), Inches(3), Inches(1.8))
    add_text(slide8, channel, left + Inches(0.15), Inches(4.15), Inches(2.7), Inches(0.35), size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide8, desc, left + Inches(0.15), Inches(4.55), Inches(2.7), Inches(1.1), size=10, color=TEXT_GRAY)

add_slide_number(slide8, 8, TOTAL_SLIDES)

# =====================================================
# SLIDE 9: TRACTION & VALIDATION
# =====================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9)

add_text(slide9, "Traction & Validation", Inches(0.5), Inches(0.3), Inches(5), Inches(0.6), size=36, color=TEXT_WHITE, bold=True)

metrics = [
    ("1,000", "Paid Users", "Year 1 Target", ACCENT_GREEN),
    ("10", "College Partners", "B2B Pipeline", ACCENT_BLUE),
    ("Rs 70L", "Year 1 Revenue", "Projected", ACCENT_ORANGE),
    ("4.8/5", "User Rating", "Beta feedback", ACCENT_PURPLE)
]

for i, (num, label, sub, color) in enumerate(metrics):
    left = Inches(0.5 + i * 3.2)
    add_card(slide9, left, Inches(1.0), Inches(3), Inches(2.0), border_color=color)
    add_text(slide9, num, left + Inches(0.15), Inches(1.2), Inches(2.7), Inches(0.6), size=36, color=color, bold=True)
    add_text(slide9, label, left + Inches(0.15), Inches(1.85), Inches(2.7), Inches(0.35), size=14, color=TEXT_WHITE, bold=True)
    add_text(slide9, sub, left + Inches(0.15), Inches(2.25), Inches(2.7), Inches(0.3), size=11, color=TEXT_GRAY)

# Milestones
add_text(slide9, "Roadmap & Milestones", Inches(0.5), Inches(3.3), Inches(4), Inches(0.4), size=18, color=TEXT_WHITE, bold=True)

milestones = [
    ("Q1 2025", "Platform Launch", "MVP with core features\nAI code generation\nBasic documentation", "Completed"),
    ("Q2 2025", "Documentation Suite", "Full SRS generation\nAll UML diagrams\nPPT & Viva Q&A", "In Progress"),
    ("Q3 2025", "College Pilots", "5 pilot institutions\nB2B dashboard\nUsage analytics", "Planned"),
    ("Q4 2025", "Scale", "20 colleges\n1,000 B2C users\nRs 70L revenue", "Target")
]

for i, (date, title, desc, status) in enumerate(milestones):
    left = Inches(0.5 + i * 3.2)

    dot = slide9.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.2), Inches(3.9), Inches(0.2), Inches(0.2))
    color = ACCENT_GREEN if status == "Completed" else (ACCENT_ORANGE if status == "In Progress" else TEXT_GRAY)
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

    add_text(slide9, date, left, Inches(4.2), Inches(3), Inches(0.3), size=11, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide9, title, left, Inches(4.5), Inches(3), Inches(0.35), size=13, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide9, desc, left, Inches(4.9), Inches(3), Inches(1.0), size=10, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide9, status, left, Inches(5.95), Inches(3), Inches(0.25), size=9, color=color, bold=True, align=PP_ALIGN.CENTER)

line = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.97), Inches(11.7), Inches(0.03))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(60, 70, 90)
line.line.fill.background()

add_slide_number(slide9, 9, TOTAL_SLIDES)

# =====================================================
# SLIDE 10: COMPETITIVE LANDSCAPE
# =====================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10)

add_text(slide10, "Competitive Advantage", Inches(0.5), Inches(0.3), Inches(5), Inches(0.6), size=36, color=TEXT_WHITE, bold=True)

headers = ["Feature", "BharatBuild", "Bolt.new", "Cursor", "Replit"]
col_widths = [3.2, 2.5, 2.3, 2.3, 2.3]

for i, (header, width) in enumerate(zip(headers, col_widths)):
    left = 0.5 + sum(col_widths[:i])
    header_box = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(0.95), Inches(width), Inches(0.45))
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = PRIMARY if i == 1 else BG_CARD_LIGHT
    header_box.line.fill.background()
    add_text(slide10, header, Inches(left), Inches(1.02), Inches(width), Inches(0.35), size=12, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)

table_data = [
    ("Pricing", "Rs 4,999 one-time", "Rs 1,700/mo", "Rs 1,700/mo", "Rs 600/mo"),
    ("Academic Docs", "Complete Package", "None", "None", "None"),
    ("SRS Generation", "Auto-generated", "Manual", "Manual", "Manual"),
    ("UML Diagrams", "All 5 types", "None", "None", "None"),
    ("Viva Q&A", "50+ questions", "None", "None", "None"),
    ("Indian Payments", "UPI/Paytm/PhonePe", "Card only", "Card only", "Card only"),
    ("Server Location", "India (Mumbai)", "USA", "USA", "USA"),
    ("Latency", "<50ms", "200-400ms", "200-400ms", "200-400ms"),
]

for row_idx, row_data in enumerate(table_data):
    top = Inches(1.45 + row_idx * 0.5)
    bg = BG_CARD if row_idx % 2 == 0 else BG_CARD_LIGHT

    for col_idx, (cell, width) in enumerate(zip(row_data, col_widths)):
        left = 0.5 + sum(col_widths[:col_idx])
        cell_box = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), top, Inches(width), Inches(0.48))
        cell_box.fill.solid()
        cell_box.fill.fore_color.rgb = bg
        cell_box.line.color.rgb = RGBColor(50, 60, 80)
        cell_box.line.width = Pt(0.5)

        if col_idx == 1:
            text_color = ACCENT_GREEN
        elif col_idx > 1 and ("None" in cell or "Card only" in cell or "Manual" in cell or "USA" in cell or "200-400ms" in cell):
            text_color = ACCENT_RED
        elif col_idx == 0:
            text_color = TEXT_WHITE
        else:
            text_color = TEXT_LIGHT

        add_text(slide10, cell, Inches(left), top + Inches(0.1), Inches(width), Inches(0.3), size=11, color=text_color, align=PP_ALIGN.CENTER)

add_text(slide10, "Why We Win", Inches(0.5), Inches(5.6), Inches(4), Inches(0.4), size=16, color=TEXT_WHITE, bold=True)

differentiators = [
    "ONLY platform with complete academic documentation (SRS, UML, Report, PPT, Viva)",
    "One-time Rs 4,999 vs competitors' Rs 1,700/month subscription model",
    "Indian servers = 4x faster response time, local payment methods",
    "Built specifically for Indian engineering curriculum requirements"
]

for i, diff in enumerate(differentiators):
    add_text(slide10, f"* {diff}", Inches(0.5), Inches(6.0 + i * 0.35), Inches(12), Inches(0.35), size=11, color=ACCENT_GREEN)

add_slide_number(slide10, 10, TOTAL_SLIDES)

# =====================================================
# SLIDE 11: TECHNOLOGY STACK
# =====================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide11)

add_text(slide11, "Technology Stack", Inches(0.5), Inches(0.3), Inches(5), Inches(0.6), size=36, color=TEXT_WHITE, bold=True)
add_text(slide11, "Built for scale, security, and speed", Inches(0.5), Inches(0.9), Inches(12), Inches(0.4), size=16, color=TEXT_GRAY)

tech_categories = [
    ("Frontend", ["Next.js 14", "React 18", "TypeScript", "Tailwind CSS", "Monaco Editor"], PRIMARY),
    ("Backend", ["FastAPI (Python)", "Node.js", "PostgreSQL", "Redis Cache", "Docker"], ACCENT_GREEN),
    ("AI/ML", ["Claude API", "GPT-4 API", "Custom prompts", "Code analysis", "Doc generation"], ACCENT_PURPLE),
    ("Infrastructure", ["AWS Mumbai", "Kubernetes", "CloudFront CDN", "Auto-scaling", "99.9% uptime"], ACCENT_ORANGE)
]

for i, (category, techs, color) in enumerate(tech_categories):
    left = Inches(0.5 + i * 3.2)
    add_card(slide11, left, Inches(1.4), Inches(3), Inches(3.8), border_color=color)
    add_text(slide11, category, left + Inches(0.15), Inches(1.55), Inches(2.7), Inches(0.4), size=16, color=color, bold=True)

    for j, tech in enumerate(techs):
        add_text(slide11, f"* {tech}", left + Inches(0.15), Inches(2.05 + j * 0.45), Inches(2.7), Inches(0.4), size=12, color=TEXT_GRAY)

# Security & Compliance
add_text(slide11, "Security & Compliance", Inches(0.5), Inches(5.4), Inches(4), Inches(0.4), size=16, color=TEXT_WHITE, bold=True)

security = ["SSL/TLS encryption", "SOC 2 compliance (planned)", "GDPR compliant", "No code stored permanently", "Secure sandboxed execution"]
for i, item in enumerate(security):
    left = Inches(0.5 + i * 2.5)
    add_text(slide11, f"* {item}", left, Inches(5.85), Inches(2.4), Inches(0.7), size=11, color=ACCENT_GREEN)

add_slide_number(slide11, 11, TOTAL_SLIDES)

# =====================================================
# SLIDE 12: FINANCIAL PROJECTIONS
# =====================================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide12)

add_text(slide12, "Financial Projections", Inches(0.5), Inches(0.3), Inches(5), Inches(0.6), size=36, color=TEXT_WHITE, bold=True)
add_text(slide12, "Based on Rs 4,999/project pricing", Inches(0.5), Inches(0.85), Inches(6), Inches(0.35), size=14, color=PRIMARY)

years = [
    ("Year 1", "Rs 70L", "1,000 paid users x Rs 4,999\n= Rs 50L (B2C)", "10 colleges x Rs 2L\n= Rs 20L (B2B)", ACCENT_BLUE),
    ("Year 2", "Rs 3.7 Cr", "5,000 paid users x Rs 4,999\n= Rs 2.5 Cr (B2C)", "40 colleges x Rs 3L\n= Rs 1.2 Cr (B2B)", ACCENT_PURPLE),
    ("Year 3", "Rs 13.5 Cr", "20,000 paid users x Rs 4,999\n= Rs 10 Cr (B2C)", "100 colleges x Rs 3.5L\n= Rs 3.5 Cr (B2B)", ACCENT_GREEN)
]

for i, (year, revenue, b2c, b2b, color) in enumerate(years):
    left = Inches(0.5 + i * 4.2)
    add_card(slide12, left, Inches(1.3), Inches(4), Inches(3.5), border_color=color)

    add_text(slide12, year, left + Inches(0.2), Inches(1.45), Inches(3.6), Inches(0.4), size=16, color=TEXT_WHITE, bold=True)
    add_text(slide12, revenue, left + Inches(0.2), Inches(1.85), Inches(3.6), Inches(0.5), size=32, color=color, bold=True)

    add_text(slide12, "B2C Revenue", left + Inches(0.2), Inches(2.5), Inches(3.6), Inches(0.3), size=11, color=ACCENT_GREEN, bold=True)
    add_text(slide12, b2c, left + Inches(0.2), Inches(2.8), Inches(3.6), Inches(0.7), size=10, color=TEXT_GRAY)

    add_text(slide12, "B2B Revenue", left + Inches(0.2), Inches(3.6), Inches(3.6), Inches(0.3), size=11, color=ACCENT_BLUE, bold=True)
    add_text(slide12, b2b, left + Inches(0.2), Inches(3.9), Inches(3.6), Inches(0.7), size=10, color=TEXT_GRAY)

# Key metrics
add_text(slide12, "Key Financial Metrics", Inches(0.5), Inches(5.0), Inches(4), Inches(0.35), size=14, color=TEXT_WHITE, bold=True)

fin_metrics = [
    ("Gross Margin", "70%"),
    ("Break-even", "Month 18"),
    ("Burn Rate", "Rs 15L/mo"),
    ("Runway", "20 months")
]

for i, (label, value) in enumerate(fin_metrics):
    left = Inches(0.5 + i * 3.2)
    add_card(slide12, left, Inches(5.4), Inches(3), Inches(0.9))
    add_text(slide12, label, left + Inches(0.1), Inches(5.5), Inches(2.8), Inches(0.25), size=10, color=TEXT_GRAY)
    add_text(slide12, value, left + Inches(0.1), Inches(5.75), Inches(2.8), Inches(0.4), size=18, color=ACCENT_GREEN, bold=True)

add_slide_number(slide12, 12, TOTAL_SLIDES)

# =====================================================
# SLIDE 13: THE ASK
# =====================================================
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide13)

add_text(slide13, "The Ask", Inches(0.5), Inches(0.3), Inches(5), Inches(0.6), size=36, color=TEXT_WHITE, bold=True)

funding_box = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(1.0), Inches(6.3), Inches(1.3))
funding_box.fill.solid()
funding_box.fill.fore_color.rgb = BG_CARD
funding_box.line.color.rgb = ACCENT_GREEN
funding_box.line.width = Pt(3)

add_text(slide13, "Seed Round: Rs 2.5 - 3 Crore", Inches(3.5), Inches(1.25), Inches(6.3), Inches(0.6), size=32, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide13, "18-20 months runway to break-even", Inches(3.5), Inches(1.85), Inches(6.3), Inches(0.35), size=14, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide13, "Use of Funds", Inches(0.5), Inches(2.6), Inches(4), Inches(0.4), size=18, color=TEXT_WHITE, bold=True)

funds = [
    ("Product & Engineering", "45%", "Rs 1.1 Cr", "AI model improvements, new features, mobile app, infrastructure scaling", ACCENT_BLUE),
    ("Sales & Marketing", "30%", "Rs 75L", "Sales team (4 FTEs), digital marketing, college outreach, ambassador program", ACCENT_GREEN),
    ("Operations", "15%", "Rs 38L", "Customer support, content creation, admin, legal & compliance", ACCENT_PURPLE),
    ("Buffer", "10%", "Rs 25L", "Contingency, working capital, unexpected opportunities", ACCENT_ORANGE)
]

for i, (category, percent, amount, desc, color) in enumerate(funds):
    top = Inches(3.1 + i * 1.0)

    bar_bg = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top + Inches(0.35), Inches(7), Inches(0.25))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = BG_CARD_LIGHT
    bar_bg.line.fill.background()

    width = float(percent.replace('%', '')) / 100 * 7
    bar_fill = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top + Inches(0.35), Inches(width), Inches(0.25))
    bar_fill.fill.solid()
    bar_fill.fill.fore_color.rgb = color
    bar_fill.line.fill.background()

    add_text(slide13, f"{category} ({percent})", Inches(0.5), top, Inches(4), Inches(0.35), size=13, color=TEXT_WHITE, bold=True)
    add_text(slide13, amount, Inches(4.5), top, Inches(1.5), Inches(0.35), size=13, color=color, bold=True)
    add_text(slide13, desc, Inches(7.8), top + Inches(0.05), Inches(5), Inches(0.55), size=10, color=TEXT_GRAY)

add_slide_number(slide13, 13, TOTAL_SLIDES)

# =====================================================
# SLIDE 14: TEAM
# =====================================================
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide14)

add_text(slide14, "Team", Inches(0.5), Inches(0.3), Inches(5), Inches(0.6), size=36, color=TEXT_WHITE, bold=True)
add_text(slide14, "Experienced founders with deep EdTech & AI expertise", Inches(0.5), Inches(0.9), Inches(12), Inches(0.4), size=16, color=PRIMARY)

team_members = [
    ("Founder & CEO", "[Name]", "10+ years in EdTech\nEx-[Company]\n[IIT/NIT] Graduate\nBuilt products used by 1M+ students"),
    ("Co-Founder & CTO", "[Name]", "8+ years in AI/ML\nEx-[Tech Company]\n[IIT/NIT] Graduate\nLed teams of 20+ engineers"),
    ("Head of Product", "[Name]", "7+ years in Product\nEx-[Startup]\n[MBA/Engineering]\nLaunched 5+ successful products"),
    ("Head of Sales", "[Name]", "6+ years in B2B Sales\nEx-[EdTech Company]\nManaged Rs 10Cr+ revenue\nStrong college network")
]

for i, (role, name, background) in enumerate(team_members):
    left = Inches(0.5 + i * 3.2)
    add_card(slide14, left, Inches(1.4), Inches(3), Inches(3.5))

    avatar = slide14.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.9), Inches(1.6), Inches(1.2), Inches(1.2))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = BG_CARD_LIGHT
    avatar.line.color.rgb = PRIMARY
    avatar.line.width = Pt(2)

    add_text(slide14, name, left, Inches(2.95), Inches(3), Inches(0.4), size=16, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide14, role, left, Inches(3.35), Inches(3), Inches(0.35), size=11, color=PRIMARY, align=PP_ALIGN.CENTER)
    add_text(slide14, background, left + Inches(0.15), Inches(3.75), Inches(2.7), Inches(1.0), size=10, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide14, "Advisors: Industry veterans from [Top EdTech Companies], [IIT Professors], [Successful Founders]",
         Inches(0.5), Inches(5.2), Inches(12), Inches(0.4), size=12, color=TEXT_GRAY)

add_text(slide14, "Hiring Plan: Year 1: 12 FTEs | Year 2: 25 FTEs | Focus: Engineering (40%), Sales (30%), Product (20%), Ops (10%)",
         Inches(0.5), Inches(5.6), Inches(12), Inches(0.4), size=12, color=TEXT_GRAY)

add_slide_number(slide14, 14, TOTAL_SLIDES)

# =====================================================
# SLIDE 15: WHY NOW?
# =====================================================
slide15 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide15)

add_text(slide15, "Why Now?", Inches(0.5), Inches(0.3), Inches(5), Inches(0.6), size=36, color=TEXT_WHITE, bold=True)
add_text(slide15, "Perfect timing for disruption in Indian EdTech", Inches(0.5), Inches(0.9), Inches(12), Inches(0.4), size=16, color=TEXT_GRAY)

why_now = [
    ("AI Revolution", "GPT-4, Claude, and other LLMs have made AI code generation viable. Quality is now production-ready. Cost has dropped 10x in 2 years.", PRIMARY, "2023-24 saw AI tools go mainstream"),
    ("NEP 2020 Implementation", "National Education Policy mandates practical, project-based learning. Colleges are actively seeking tools to help students build real projects.", ACCENT_GREEN, "Policy tailwind for our solution"),
    ("UPI Penetration", "500M+ Indians now use UPI. Students in Tier 2/3 cities can pay digitally. No more credit card barrier for Indian products.", ACCENT_ORANGE, "Payment infrastructure ready"),
    ("Post-COVID Digital Adoption", "Students and colleges are now comfortable with online tools. EdTech acceptance at all-time high. Remote learning is normalized.", ACCENT_PURPLE, "Behavioral shift in our favor")
]

for i, (title, desc, color, insight) in enumerate(why_now):
    col = i % 2
    row = i // 2
    left = Inches(0.5 + col * 6.4)
    top = Inches(1.5 + row * 2.8)

    add_card(slide15, left, top, Inches(6.2), Inches(2.6), border_color=color)
    add_text(slide15, title, left + Inches(0.25), top + Inches(0.2), Inches(5.7), Inches(0.45), size=18, color=color, bold=True)
    add_text(slide15, desc, left + Inches(0.25), top + Inches(0.7), Inches(5.7), Inches(1.2), size=12, color=TEXT_GRAY)
    add_text(slide15, insight, left + Inches(0.25), top + Inches(2.0), Inches(5.7), Inches(0.4), size=11, color=color, bold=True)

add_slide_number(slide15, 15, TOTAL_SLIDES)

# =====================================================
# SLIDE 16: VISION & CONTACT
# =====================================================
slide16 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide16)

add_text(slide16, "Our Vision", Inches(0.5), Inches(1.2), prs.slide_width - Inches(1), Inches(0.5), size=20, color=PRIMARY, align=PP_ALIGN.CENTER)

add_text(slide16, "Democratize AI-powered development\nfor every Indian student", Inches(0.5), Inches(1.8), prs.slide_width - Inches(1), Inches(1.2),
         size=40, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide16, "From Tier 3 towns to metro cities, every engineering student deserves\naccess to world-class AI tools at Indian prices.", Inches(0.5), Inches(3.2), prs.slide_width - Inches(1), Inches(0.8),
         size=16, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# Summary points
points = [
    ("Rs 2,195 Cr", "Market Opportunity"),
    ("Rs 4,999", "Per Project"),
    ("Rs 13.5 Cr", "Year 3 Revenue"),
    ("Rs 2.5-3 Cr", "Seed Ask")
]

for i, (value, label) in enumerate(points):
    left = Inches(1.5 + i * 2.8)
    add_card(slide16, left, Inches(4.2), Inches(2.5), Inches(1.3))
    add_text(slide16, value, left, Inches(4.35), Inches(2.5), Inches(0.5), size=18, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide16, label, left, Inches(4.85), Inches(2.5), Inches(0.6), size=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

add_card(slide16, Inches(4), Inches(5.8), Inches(5.3), Inches(1.2))
add_text(slide16, "Let's Build India's Future Together", Inches(4), Inches(5.95), Inches(5.3), Inches(0.45), size=18, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide16, "www.bharatbuild.ai | contact@bharatbuild.ai", Inches(4), Inches(6.45), Inches(5.3), Inches(0.4), size=14, color=PRIMARY, align=PP_ALIGN.CENTER)

add_slide_number(slide16, 16, TOTAL_SLIDES)

# =====================================================
# SAVE PRESENTATION
# =====================================================
output_path = os.path.join(os.path.dirname(__file__), "BharatBuild_Investor_Deck_v2.pptx")
prs.save(output_path)
print(f"\nInvestor Pitch Deck saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print("\nSlide Contents:")
print("  1. Title / Cover")
print("  2. The Problem (4 key challenges)")
print("  3. Our Solution (6 features)")
print("  4. Product Screenshots (5 pages)")
print("  5. How It Works (4 steps)")
print("  6. Market Opportunity (TAM/SAM/SOM)")
print("  7. Business Model (B2C + B2B)")
print("  8. Go-to-Market Strategy")
print("  9. Traction & Roadmap")
print(" 10. Competitive Landscape")
print(" 11. Technology Stack")
print(" 12. Financial Projections (3 years)")
print(" 13. The Ask (Use of Funds)")
print(" 14. Team")
print(" 15. Why Now?")
print(" 16. Vision & Contact")
